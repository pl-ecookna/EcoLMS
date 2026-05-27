from __future__ import annotations

import base64
import json
import os
import socket
import ssl
import subprocess
import tempfile
import time
import uuid
from io import BytesIO
from dataclasses import asdict, dataclass
from datetime import datetime, timezone
from typing import Any
from urllib import error as urlerror
from urllib import request as urlrequest
from urllib.parse import urlparse, unquote

import boto3
import certifi
import psycopg
from botocore.config import Config
from docx import Document
from pptx import Presentation
from pypdf import PdfReader
from psycopg.rows import dict_row
from striprtf.striprtf import rtf_to_text

from worker.prompts import DEFAULT_PROMPTS, PromptDefinition, prompt_bundle_for_stage


@dataclass(slots=True)
class WorkerConfig:
    api_base_url: str = "http://localhost:3101"
    transcription_service_url: str = "http://localhost:3002"
    postgres_url: str = "postgresql://postgres:password@localhost:5432/ecolms"
    redis_url: str = "redis://localhost:6379"
    s3_endpoint: str = "https://s3.example.invalid"
    s3_bucket: str = "ecolms"
    s3_region: str = "ru-1"
    s3_access_key_id: str = ""
    s3_secret_access_key: str = ""
    openai_api_key: str = ""
    openrouter_api_key: str = ""
    llm_primary_provider: str = "openai"
    openai_model: str = "gpt-4.1-mini"
    openrouter_model: str = "openai/gpt-4.1-mini"
    openrouter_base_url: str = "https://openrouter.ai/api/v1/chat/completions"
    llm_timeout_seconds: float = 120.0
    meeting_transcription_provider: str = "assemblyai"
    job_queue_key: str = "ecolms:processing-jobs"
    meeting_job_queue_key: str = "ecolms:meeting-jobs"
    assemblyai_api_key: str = ""
    assemblyai_base_url: str = "https://api.eu.assemblyai.com"
    assemblyai_speech_models: tuple[str, ...] = ("universal",)
    assemblyai_language_code: str = "ru"
    assemblyai_poll_interval_seconds: float = 5.0
    assemblyai_timeout_seconds: float = 1800.0
    assemblyai_audio_url_expires_seconds: int = 1800
    meeting_audio_prep_timeout_seconds: float = 900.0
    salutespeech_auth_key: str = ""
    salutespeech_oauth_url: str = ""
    salutespeech_rest_url: str = "https://smartspeech.sber.ru/rest/v1"
    salutespeech_upload_url: str = "https://smartspeech.sber.ru/rest/v1/data:upload"
    salutespeech_recognize_url: str = "https://smartspeech.sber.ru/rest/v1/speech:async_recognize"
    salutespeech_task_url: str = "https://smartspeech.sber.ru/rest/v1/task:get"
    salutespeech_download_url: str = "https://smartspeech.sber.ru/rest/v1/data:download"
    salutespeech_scope: str = "SALUTE_SPEECH_PERS"
    salutespeech_model: str = "general"
    salutespeech_language: str = "ru-RU"
    salutespeech_ca_cert_path: str = ""
    salutespeech_ssl_no_verify: bool = False
    salutespeech_poll_interval_seconds: float = 5.0
    salutespeech_timeout_seconds: float = 1800.0


INTERNAL_POSTGRES_URL = (
    "postgresql://postgres:postgres@postgres:5432/ecolms"
)
EXTERNAL_POSTGRES_URL = (
    "postgresql://postgres:postgres@localhost:5434/ecolms"
)
INTERNAL_REDIS_URL = "redis://redis:6379"
EXTERNAL_REDIS_URL = "redis://localhost:6381"
_LLM_SSL_CONTEXT: ssl.SSLContext | None = None


def llm_ssl_context() -> ssl.SSLContext:
    global _LLM_SSL_CONTEXT
    if _LLM_SSL_CONTEXT is None:
        _LLM_SSL_CONTEXT = ssl.create_default_context(cafile=certifi.where())
    return _LLM_SSL_CONTEXT


def normalize_service_url(value: str | None, *, default_prod: str, default_dev: str) -> str:
    if not value:
        return default_prod if os.getenv("NODE_ENV") == "production" else default_dev

    if any(
        marker in value
        for marker in (
            "api:3001",
            "transcription-service:3002",
            "localhost:3001",
            "localhost:3002",
            "127.0.0.1:3001",
            "127.0.0.1:3002",
        )
    ):
        return default_prod if os.getenv("NODE_ENV") == "production" else default_dev

    return value


def first_defined_env(*names: str) -> str:
    for name in names:
        value = os.getenv(name, "").strip()
        if value:
            return value
    return ""


def normalize_basic_auth_value(value: str) -> str:
    candidate = value.strip()
    if not candidate:
        return ""
    if candidate.lower().startswith("basic "):
        candidate = candidate[6:].strip()
    if ":" in candidate:
        return base64.b64encode(candidate.encode("utf-8")).decode("ascii")
    return candidate


def resolve_salutespeech_auth_key() -> str:
    auth_key = normalize_basic_auth_value(
        first_defined_env("SALUTESPEECH_AUTH_KEY", "SBER_AUTH_KEY")
    )
    if auth_key:
        return auth_key

    client_id = first_defined_env("SALUTESPEECH_CLIENT_ID", "SBER_CLIENT_ID")
    client_secret = first_defined_env(
        "SALUTESPEECH_CLIENT_SECRET",
        "SBER_CLIENT_SECRET",
    )
    if client_id and client_secret:
        raw = f"{client_id}:{client_secret}".encode("utf-8")
        return base64.b64encode(raw).decode("ascii")
    return ""


def parse_csv_env(value: str, fallback: tuple[str, ...]) -> tuple[str, ...]:
    items = [item.strip() for item in value.split(",") if item.strip()]
    return tuple(items) if items else fallback


def normalize_meeting_transcription_provider(value: str) -> str:
    provider = value.strip().lower() or "assemblyai"
    if provider not in {"assemblyai", "salutespeech"}:
        raise RuntimeError(
            "MEETING_TRANSCRIPTION_PROVIDER должен быть assemblyai или salutespeech."
        )
    return provider


def load_config() -> WorkerConfig:
    salutespeech_auth_key = resolve_salutespeech_auth_key()
    salutespeech_rest_url = first_defined_env(
        "SALUTESPEECH_REST_URL",
        "SBER_REST_URL",
    ) or os.getenv(
        "SALUTESPEECH_REST_URL",
        "https://smartspeech.sber.ru/rest/v1",
    ).strip()
    meeting_transcription_provider = normalize_meeting_transcription_provider(
        os.getenv("MEETING_TRANSCRIPTION_PROVIDER", "assemblyai")
    )
    return WorkerConfig(
        api_base_url=normalize_service_url(
            os.getenv("API_BASE_URL"),
            default_prod="http://app-calculate-open-source-alarm-cob2f6:3001",
            default_dev="http://localhost:3101",
        ),
        transcription_service_url=normalize_service_url(
            os.getenv("TRANSCRIPTION_SERVICE_URL"),
            default_prod="http://app-copy-bluetooth-matrix-b869ye:3002",
            default_dev="http://localhost:3002",
        ),
        postgres_url=normalize_service_url(
            os.getenv("POSTGRES_URL"),
            default_prod=INTERNAL_POSTGRES_URL,
            default_dev=EXTERNAL_POSTGRES_URL,
        ),
        redis_url=normalize_service_url(
            os.getenv("REDIS_URL"),
            default_prod=INTERNAL_REDIS_URL,
            default_dev=EXTERNAL_REDIS_URL,
        ),
        s3_endpoint=os.getenv("S3_ENDPOINT", "https://s3.example.invalid"),
        s3_bucket=os.getenv("S3_BUCKET", "ecolms"),
        s3_region=os.getenv("S3_REGION", "ru-1"),
        s3_access_key_id=os.getenv("S3_ACCESS_KEY_ID", ""),
        s3_secret_access_key=os.getenv("S3_SECRET_ACCESS_KEY", ""),
        openai_api_key=os.getenv("OPENAI_API_KEY", ""),
        openrouter_api_key=os.getenv("OPENROUTER_API_KEY", ""),
        llm_primary_provider=os.getenv("LLM_PRIMARY_PROVIDER", "openai").strip().lower(),
        openai_model=os.getenv("OPENAI_MODEL", "gpt-4.1-mini"),
        openrouter_model=os.getenv("OPENROUTER_MODEL", "openai/gpt-4.1-mini"),
        openrouter_base_url=os.getenv(
            "OPENROUTER_BASE_URL", "https://openrouter.ai/api/v1/chat/completions"
        ),
        llm_timeout_seconds=float(os.getenv("LLM_TIMEOUT_SECONDS", "120")),
        meeting_transcription_provider=meeting_transcription_provider,
        job_queue_key=os.getenv("WORKER_JOB_QUEUE_KEY", "ecolms:processing-jobs"),
        meeting_job_queue_key=os.getenv("WORKER_MEETING_JOB_QUEUE_KEY", "ecolms:meeting-jobs"),
        assemblyai_api_key=os.getenv("ASSEMBLYAI_API_KEY", "").strip(),
        assemblyai_base_url=os.getenv(
            "ASSEMBLYAI_BASE_URL", "https://api.eu.assemblyai.com"
        ).strip().rstrip("/"),
        assemblyai_speech_models=parse_csv_env(
            os.getenv("ASSEMBLYAI_SPEECH_MODELS", "universal"),
            ("universal",),
        ),
        assemblyai_language_code=os.getenv("ASSEMBLYAI_LANGUAGE_CODE", "ru").strip() or "ru",
        assemblyai_poll_interval_seconds=float(
            os.getenv("ASSEMBLYAI_POLL_INTERVAL_SECONDS", "5")
        ),
        assemblyai_timeout_seconds=float(os.getenv("ASSEMBLYAI_TIMEOUT_SECONDS", "1800")),
        assemblyai_audio_url_expires_seconds=int(
            os.getenv("ASSEMBLYAI_AUDIO_URL_EXPIRES_SECONDS", "1800")
        ),
        meeting_audio_prep_timeout_seconds=float(
            os.getenv("MEETING_AUDIO_PREP_TIMEOUT_SECONDS", "900")
        ),
        salutespeech_auth_key=salutespeech_auth_key,
        salutespeech_oauth_url=first_defined_env(
            "SALUTESPEECH_OAUTH_URL",
            "SBER_OAUTH_URL",
        ),
        salutespeech_rest_url=salutespeech_rest_url,
        salutespeech_upload_url=os.getenv(
            "SALUTESPEECH_UPLOAD_URL",
            f"{salutespeech_rest_url.rstrip('/')}/data:upload",
        ).strip(),
        salutespeech_recognize_url=os.getenv(
            "SALUTESPEECH_RECOGNIZE_URL",
            f"{salutespeech_rest_url.rstrip('/')}/speech:async_recognize",
        ).strip(),
        salutespeech_task_url=os.getenv(
            "SALUTESPEECH_TASK_URL",
            f"{salutespeech_rest_url.rstrip('/')}/task:get",
        ).strip(),
        salutespeech_download_url=os.getenv(
            "SALUTESPEECH_DOWNLOAD_URL",
            f"{salutespeech_rest_url.rstrip('/')}/data:download",
        ).strip(),
        salutespeech_scope=first_defined_env(
            "SALUTESPEECH_SCOPE",
            "SBER_SCOPE",
        ) or "SALUTE_SPEECH_PERS",
        salutespeech_model=os.getenv("SALUTESPEECH_MODEL", "general").strip() or "general",
        salutespeech_language=os.getenv("SALUTESPEECH_LANGUAGE", "ru-RU").strip() or "ru-RU",
        salutespeech_ca_cert_path=first_defined_env(
            "SALUTESPEECH_CA_CERT_PATH",
            "SBER_CA_CERT_PATH",
        ),
        salutespeech_ssl_no_verify=os.getenv("SALUTESPEECH_SSL_NO_VERIFY", "false").strip().lower()
        in {"1", "true", "yes", "on"},
        salutespeech_poll_interval_seconds=float(
            os.getenv("SALUTESPEECH_POLL_INTERVAL_SECONDS", "5")
        ),
        salutespeech_timeout_seconds=float(
            os.getenv("SALUTESPEECH_TIMEOUT_SECONDS", "1800")
        ),
    )


def utc_now() -> str:
    return datetime.now(timezone.utc).isoformat()


def encode_resp(parts: list[str]) -> bytes:
    payload = [f"*{len(parts)}\r\n"]
    for part in parts:
        encoded = part.encode("utf-8")
        payload.append(f"${len(encoded)}\r\n")
        payload.append(part)
        payload.append("\r\n")
    return "".join(payload).encode("utf-8")


def parse_resp(buffer: bytes) -> tuple[Any, bytes] | None:
    if not buffer:
        return None

    prefix = buffer[:1]
    if prefix == b"+":
        end = buffer.find(b"\r\n")
        if end == -1:
            return None
        return buffer[1:end].decode("utf-8"), buffer[end + 2 :]

    if prefix == b":":
        end = buffer.find(b"\r\n")
        if end == -1:
            return None
        return int(buffer[1:end]), buffer[end + 2 :]

    if prefix == b"$":
        end = buffer.find(b"\r\n")
        if end == -1:
            return None
        length = int(buffer[1:end])
        if length == -1:
            return None, buffer[end + 2 :]
        start = end + 2
        stop = start + length
        if len(buffer) < stop + 2:
            return None
        return buffer[start:stop].decode("utf-8"), buffer[stop + 2 :]

    if prefix == b"*":
        end = buffer.find(b"\r\n")
        if end == -1:
            return None
        count = int(buffer[1:end])
        if count == -1:
            return None, buffer[end + 2 :]

        cursor = end + 2
        items: list[Any] = []
        for _ in range(count):
            decoded = parse_resp(buffer[cursor:])
            if decoded is None:
                return None
            value, rest = decoded
            items.append(value)
            cursor = len(buffer) - len(rest)
        return items, buffer[cursor:]

    if prefix == b"-":
        end = buffer.find(b"\r\n")
        if end == -1:
            return None
        raise RuntimeError(buffer[1:end].decode("utf-8"))

    return None


def redis_url_parts(redis_url: str) -> tuple[str, int, str | None, str | None, int]:
    parsed = urlparse(redis_url)
    host = parsed.hostname or "localhost"
    port = parsed.port or 6379
    username = unquote(parsed.username) if parsed.username else None
    password = unquote(parsed.password) if parsed.password else None
    db = 0
    if parsed.path and parsed.path != "/":
        try:
            db = int(parsed.path.lstrip("/"))
        except ValueError:
            db = 0
    return host, port, username, password, db


def redis_command(redis_url: str, *parts: str, timeout: float = 6.0) -> Any:
    host, port, username, password, db = redis_url_parts(redis_url)
    with socket.create_connection((host, port), timeout=timeout) as sock:
        sock.settimeout(timeout)
        if password:
            auth_parts = ["AUTH"]
            if username:
                auth_parts.extend([username, password])
            else:
                auth_parts.append(password)
            sock.sendall(encode_resp(auth_parts))
            auth_buffer = b""
            while True:
                chunk = sock.recv(4096)
                if not chunk:
                    raise RuntimeError("Redis connection closed during AUTH")
                auth_buffer += chunk
                decoded = parse_resp(auth_buffer)
                if decoded is not None:
                    break

        if db > 0:
            sock.sendall(encode_resp(["SELECT", str(db)]))
            select_buffer = b""
            while True:
                chunk = sock.recv(4096)
                if not chunk:
                    raise RuntimeError("Redis connection closed during SELECT")
                select_buffer += chunk
                decoded = parse_resp(select_buffer)
                if decoded is not None:
                    break

        sock.sendall(encode_resp(list(parts)))
        buffer = b""
        while True:
            chunk = sock.recv(4096)
            if not chunk:
                raise RuntimeError("Redis connection closed")
            buffer += chunk
            decoded = parse_resp(buffer)
            if decoded is not None:
                value, _ = decoded
                return value


def make_stage_markdown(project_name: str, stage: str, topic: str) -> str:
    if stage == "source_compiled":
        return f"# {project_name}\n\n{topic}\n"

    if stage == "course_outline":
        return (
            "# План курса\n\n"
            f"**Основа курса:** {topic}\n\n"
            "## Разделы курса\n"
            "1. Введение в продукт, систему или технологию\n"
            "2. Ключевые технические сведения и параметры\n"
            "3. Основные операции и порядок выполнения работ\n"
            "4. Контроль качества и типовые ошибки\n"
            "5. Безопасность и итоговая проверка знаний\n\n"
            "## Ключевые навыки после обучения\n"
            "- ориентироваться в исходных материалах;\n"
            "- понимать последовательность работ;\n"
            "- видеть требования к качеству и безопасности.\n"
        )

    if stage == "course_content":
        return (
            "# Обучающие материалы\n\n"
            "## Цели обучения\n"
            "- Сформировать прикладное понимание темы на основе утверждённого плана.\n"
            "- Подготовить материал для практического использования в работе.\n\n"
            "## Теоретическая часть\n"
            f"{topic}\n\n"
            "## Практическое применение\n"
            "1. Выполнить действия в порядке, описанном в исходных материалах.\n"
            "2. Проверить соблюдение технических требований и допусков.\n"
            "3. Зафиксировать контрольные точки качества.\n\n"
            "## Важные замечания\n"
            "> Если данных из исходников недостаточно для полного раскрытия раздела, это нужно явно отметить при доработке материала.\n"
        )

    return (
        "# Тест\n\n"
        "1. Какой принцип обязателен при создании учебного контента?\n"
        "   - Использовать только данные из исходных материалов ✅\n"
        "   - Добавлять сведения из внешних источников\n"
        "   - Сокращать технические требования по усмотрению автора\n\n"
        "2. Что должно быть отражено в хорошем учебном материале?\n"
        "   - Только маркетинговые преимущества\n"
        "   - Контроль качества, безопасность и порядок действий ✅\n"
        "   - Только перечень файлов проекта\n"
    )


def make_source_compiled_markdown(project_name: str, sections: list[str]) -> str:
    clean_sections = [item.strip() for item in sections if item and item.strip()]
    if not clean_sections:
        return f"# {project_name}\n\nИсточник пуст. Добавьте файлы и запустите распознавание.\n"
    return f"# {project_name}\n\n" + "\n\n---\n\n".join(clean_sections)


def extract_summary(text: str, limit: int = 1500) -> str:
    compact = " ".join(text.split())
    if not compact:
        return "Транскрипт получен, но текст пуст."
    if len(compact) <= limit:
        return compact
    return f"{compact[:limit].rstrip()}..."


def parse_json_from_text(value: str) -> dict[str, Any]:
    raw = value.strip()
    if not raw:
        raise RuntimeError("Пустой ответ LLM.")
    try:
        parsed = json.loads(raw)
    except json.JSONDecodeError:
        start = raw.find("{")
        end = raw.rfind("}")
        if start == -1 or end == -1 or end <= start:
            raise RuntimeError("LLM вернул невалидный JSON.")
        parsed = json.loads(raw[start : end + 1])
    if not isinstance(parsed, dict):
        raise RuntimeError("LLM вернул невалидную JSON-структуру.")
    return parsed


def resolve_llm_provider(config: WorkerConfig) -> str:
    primary = config.llm_primary_provider
    if primary not in ("openai", "openrouter"):
        primary = "openai"

    if primary == "openai" and not config.openai_api_key.strip():
        raise RuntimeError("Выбран LLM-провайдер openai, но OPENAI_API_KEY не задан.")
    if primary == "openrouter" and not config.openrouter_api_key.strip():
        raise RuntimeError(
            "Выбран LLM-провайдер openrouter, но OPENROUTER_API_KEY не задан."
        )

    return primary


def ensure_prompt_storage(conn: psycopg.Connection) -> None:
    conn.execute(
        """
        create table if not exists llm_prompts (
          module text not null check (module in ('lms', 'meetings')),
          prompt_key text not null,
          title text not null,
          system_prompt text not null,
          user_prompt_template text not null default '',
          created_at timestamptz not null default now(),
          updated_at timestamptz not null default now(),
          primary key (module, prompt_key)
        )
        """
    )
    conn.execute(
        """
        create index if not exists llm_prompts_module_idx
          on llm_prompts (module, prompt_key)
        """
    )

    for prompt in DEFAULT_PROMPTS.values():
        conn.execute(
            """
            insert into llm_prompts (
              module, prompt_key, title, system_prompt, user_prompt_template, created_at, updated_at
            ) values (
              %s, %s, %s, %s, %s, now(), now()
            )
            on conflict (module, prompt_key) do nothing
            """,
            (
                prompt.module,
                prompt.key,
                prompt.title,
                prompt.system_prompt,
                prompt.user_prompt_template,
            ),
        )

    default_prompt_upgrades = {
        "meeting_summary": "Верни только JSON",
        "meeting_protocol": "Верни только JSON",
        "meeting_actions": (
            "Не придумывай имена, сроки, решения и поручения, которых нет в тексте."
        ),
    }
    for prompt_key, old_system_needle in default_prompt_upgrades.items():
        prompt = DEFAULT_PROMPTS.get(("meetings", prompt_key))
        if prompt is None:
            continue
        conn.execute(
            """
            update llm_prompts
            set
              title = %s,
              system_prompt = %s,
              user_prompt_template = %s,
              updated_at = now()
            where module = %s
              and prompt_key = %s
              and position(%s in system_prompt) > 0
            """,
            (
                prompt.title,
                prompt.system_prompt,
                prompt.user_prompt_template,
                prompt.module,
                prompt.key,
                old_system_needle,
            ),
        )


def load_prompt_definition(
    conn: psycopg.Connection,
    module: str,
    prompt_key: str,
) -> PromptDefinition:
    ensure_prompt_storage(conn)
    row = conn.execute(
        """
        select module, prompt_key, title, system_prompt, user_prompt_template
        from llm_prompts
        where module = %s and prompt_key = %s
        limit 1
        """,
        (module, prompt_key),
    ).fetchone()
    if row is None:
        default_prompt = DEFAULT_PROMPTS.get((module, prompt_key))
        if default_prompt is None:
            raise RuntimeError(f"Неизвестный prompt: {module}:{prompt_key}")
        return default_prompt
    return PromptDefinition(
        module=str(row["module"]),
        key=str(row["prompt_key"]),
        title=str(row["title"]),
        system_prompt=str(row["system_prompt"]),
        user_prompt_template=str(row["user_prompt_template"] or ""),
    )


def normalize_provider_error(provider: str, raw_error: str) -> str:
    normalized = raw_error.lower()
    billing_patterns = [
        "insufficient_quota",
        "insufficient quota",
        "insufficient funds",
        "out of credits",
        "payment required",
        "billing",
        "quota exceeded",
        "not enough credits",
        "credit balance",
        "недостаточно средств",
        "недостаточно денег",
        "недостаточно кредитов",
        "исчерпан лимит",
        "квота",
        "rate limit",
        "too many requests",
        '"status":402',
        "http 402",
    ]
    if any(pattern in normalized for pattern in billing_patterns):
        if provider == "salutespeech":
            return "SaluteSpeech недоступен: недостаточно средств или исчерпана квота."
        if provider == "assemblyai":
            return "AssemblyAI недоступен: недостаточно средств, исчерпана квота или превышен лимит запросов."
        return f"{provider} недоступен: недостаточно средств или исчерпана квота."
    return raw_error


def parse_json(value: Any, fallback: dict[str, Any]) -> dict[str, Any]:
    if isinstance(value, dict):
        return value
    if isinstance(value, str):
        raw = value.strip()
        if not raw:
            return fallback
        try:
            parsed = json.loads(raw)
        except json.JSONDecodeError:
            return fallback
        if isinstance(parsed, dict):
            return parsed
    return fallback


def openai_chat_completion(
    api_key: str, model: str, messages: list[dict[str, str]], timeout_seconds: float
) -> tuple[str, str]:
    request_payload = {
        "model": model,
        "messages": messages,
        "temperature": 0.2,
        "response_format": {"type": "json_object"},
    }
    req = urlrequest.Request(
        "https://api.openai.com/v1/chat/completions",
        method="POST",
        data=json.dumps(request_payload, ensure_ascii=False).encode("utf-8"),
        headers={
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json; charset=utf-8",
        },
    )
    try:
        with urlrequest.urlopen(
            req,
            timeout=timeout_seconds,
            context=llm_ssl_context(),
        ) as response:
            body = json.loads(response.read().decode("utf-8"))
    except urlerror.HTTPError as exc:
        try:
            details = exc.read().decode("utf-8", errors="ignore").strip()
        except Exception:
            details = ""
        message = normalize_provider_error(
            "OpenAI",
            f"HTTP {exc.code}{f': {details}' if details else ''}",
        )
        raise RuntimeError(message) from exc
    content = (
        body.get("choices", [{}])[0]
        .get("message", {})
        .get("content", "")
    )
    if not isinstance(content, str) or not content.strip():
        raise RuntimeError("OpenAI вернул пустой ответ.")
    model_name = str(body.get("model") or model)
    return content, model_name


def openrouter_chat_completion(
    api_key: str,
    base_url: str,
    model: str,
    messages: list[dict[str, str]],
    timeout_seconds: float,
) -> tuple[str, str]:
    request_payload = {
        "model": model,
        "messages": messages,
        "temperature": 0.2,
        "response_format": {"type": "json_object"},
    }
    req = urlrequest.Request(
        base_url,
        method="POST",
        data=json.dumps(request_payload, ensure_ascii=False).encode("utf-8"),
        headers={
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json; charset=utf-8",
            "HTTP-Referer": "https://ecolms.local",
            "X-Title": "EcoLMS Worker",
        },
    )
    try:
        with urlrequest.urlopen(
            req,
            timeout=timeout_seconds,
            context=llm_ssl_context(),
        ) as response:
            body = json.loads(response.read().decode("utf-8"))
    except urlerror.HTTPError as exc:
        try:
            details = exc.read().decode("utf-8", errors="ignore").strip()
        except Exception:
            details = ""
        message = normalize_provider_error(
            "OpenRouter",
            f"HTTP {exc.code}{f': {details}' if details else ''}",
        )
        raise RuntimeError(message) from exc
    content = (
        body.get("choices", [{}])[0]
        .get("message", {})
        .get("content", "")
    )
    if not isinstance(content, str) or not content.strip():
        raise RuntimeError("OpenRouter вернул пустой ответ.")
    model_name = str(body.get("model") or model)
    return content, model_name


def analyze_source_with_llm(
    conn: psycopg.Connection,
    config: WorkerConfig,
    source_text: str,
    *,
    prompt_key: str = "analize_video",
    source_type: str = "видео",
) -> dict[str, Any]:
    prompt_definition = load_prompt_definition(conn, "lms", prompt_key)
    prompt = prompt_definition.system_prompt
    if not source_text.strip():
        raise RuntimeError("Невозможно вызвать LLM: исходный текст пуст.")

    user_payload = {
        "task": prompt_definition.user_prompt_template.replace("{source_type}", source_type),
        "sourceText": source_text,
    }
    messages = [
        {"role": "system", "content": prompt},
        {"role": "user", "content": json.dumps(user_payload, ensure_ascii=False)},
    ]

    provider = resolve_llm_provider(config)
    try:
        if provider == "openai":
            content, model_name = openai_chat_completion(
                config.openai_api_key,
                config.openai_model,
                messages,
                config.llm_timeout_seconds,
            )
        else:
            content, model_name = openrouter_chat_completion(
                config.openrouter_api_key,
                config.openrouter_base_url,
                config.openrouter_model,
                messages,
                config.llm_timeout_seconds,
            )

        parsed = parse_json_from_text(content)
        summary = str(parsed.get("summary") or "").strip()
        transcript = str(
            parsed.get("transcript")
            or parsed.get("sourceText")
            or parsed.get("text")
            or ""
        ).strip()
        if not summary:
            summary = extract_summary(source_text)
        if not transcript:
            transcript = source_text
        return {
            "summary": summary,
            "transcript": transcript,
            "provider": provider,
            "model": model_name,
            "attempts": [{"provider": provider, "ok": True}],
            "promptModule": prompt_definition.module,
            "promptKey": prompt_definition.key,
            "promptTitle": prompt_definition.title,
        }
    except Exception as exc:
        raise RuntimeError(f"LLM-провайдер {provider} завершился ошибкой: {exc}") from exc


def stage_input_stages(stage: str) -> list[str]:
    if stage == "course_outline":
        return ["source_compiled"]
    if stage == "course_content":
        return ["source_compiled", "course_outline"]
    if stage == "course_test":
        return ["source_compiled", "course_content"]
    return []


def build_stage_generation_input(
    conn: psycopg.Connection, project_id: str, stage: str, fallback_summary: str
) -> tuple[str, list[str]]:
    stages = stage_input_stages(stage)
    parts: list[str] = []
    used_stages: list[str] = []
    for source_stage in stages:
        markdown = load_artifact_markdown(conn, project_id, source_stage)
        if not markdown:
            continue
        parts.append(f"### {source_stage}\n{markdown}")
        used_stages.append(source_stage)
    if parts:
        return "\n\n".join(parts).strip(), used_stages
    return fallback_summary.strip(), used_stages


def generate_stage_markdown_with_llm(
    conn: psycopg.Connection,
    config: WorkerConfig,
    *,
    stage: str,
    source_text: str,
    project_name: str,
) -> dict[str, Any]:
    prompt_keys = [item.key for item in prompt_bundle_for_stage(stage)]
    prompt_key = next((key for key in prompt_keys if key.startswith("generate_")), None)
    if not prompt_key:
        raise RuntimeError(f"Для этапа {stage} не найден prompt генерации.")
    prompt_definition = load_prompt_definition(conn, "lms", prompt_key)
    prompt = prompt_definition.system_prompt
    normalized_source = source_text.strip()
    if not normalized_source:
        raise RuntimeError(f"Невозможно вызвать LLM для {stage}: пустой источник.")

    user_payload = {
        "task": prompt_definition.user_prompt_template,
        "projectName": project_name,
        "stage": stage,
        "sourceText": normalized_source[:180000],
        "outputFormat": {"type": "json", "fields": ["markdown", "shortSummary"]},
    }
    messages = [
        {"role": "system", "content": prompt},
        {"role": "user", "content": json.dumps(user_payload, ensure_ascii=False)},
    ]

    provider = resolve_llm_provider(config)
    try:
        if provider == "openai":
            content, model_name = openai_chat_completion(
                config.openai_api_key,
                config.openai_model,
                messages,
                config.llm_timeout_seconds,
            )
        else:
            content, model_name = openrouter_chat_completion(
                config.openrouter_api_key,
                config.openrouter_base_url,
                config.openrouter_model,
                messages,
                config.llm_timeout_seconds,
            )
        parsed = parse_json_from_text(content)
        markdown = str(
            parsed.get("markdown") or parsed.get("content") or parsed.get("text") or ""
        ).strip()
        short_summary = str(parsed.get("shortSummary") or "").strip()
        if not markdown:
            raise RuntimeError("LLM вернул пустое поле markdown.")
        return {
            "markdown": markdown,
            "shortSummary": short_summary,
            "provider": provider,
            "model": model_name,
            "attempts": [{"provider": provider, "ok": True}],
            "promptModule": prompt_definition.module,
            "sourceTextLength": len(normalized_source),
            "promptKey": prompt_key,
            "promptTitle": prompt_definition.title,
        }
    except Exception as exc:
        raise RuntimeError(f"LLM-провайдер {provider} завершился ошибкой: {exc}") from exc


def stage_prompt_metadata(conn: psycopg.Connection, stage: str) -> dict[str, Any]:
    prompts = [
        load_prompt_definition(conn, "lms", prompt.key)
        for prompt in prompt_bundle_for_stage(stage)
    ]
    return {
        "promptModules": [item.module for item in prompts],
        "promptKeys": [item.key for item in prompts],
        "promptTitles": [item.title for item in prompts],
        "promptTexts": [item.system_prompt for item in prompts],
        "promptUserTemplates": [item.user_prompt_template for item in prompts],
    }


STAGE_ORDER = ["source_compiled", "course_outline", "course_content", "course_test"]
MEETING_STAGE_ORDER = [
    "audio_prepared",
    "transcript_compiled",
    "meeting_summary",
    "meeting_protocol",
    "meeting_actions",
]


def next_stage(stage: str) -> str | None:
    try:
        index = STAGE_ORDER.index(stage)
    except ValueError:
        return None
    return STAGE_ORDER[index + 1] if index + 1 < len(STAGE_ORDER) else None


def next_meeting_stage(stage: str) -> str | None:
    try:
        index = MEETING_STAGE_ORDER.index(stage)
    except ValueError:
        return None
    return (
        MEETING_STAGE_ORDER[index + 1]
        if index + 1 < len(MEETING_STAGE_ORDER)
        else None
    )


def progress_for_stage(stage: str, *, completed: bool) -> int:
    if completed:
        return 100
    if stage == "source_compiled":
        return 18
    if stage == "course_outline":
        return 44
    if stage == "course_content":
        return 68
    if stage == "course_test":
        return 88
    return 0


def ensure_stage_artifact(
    conn: psycopg.Connection, job: dict[str, Any], markdown: str, metadata: dict[str, Any]
) -> None:
    conn.execute(
        """
        update artifacts
        set content_md = %s,
            content_json = %s::jsonb,
            updated_at = now()
        where project_id = %s and stage = %s and format = 'md'
        """,
        (
            markdown,
            json.dumps(
                {"stage": job["stage"], "markdown": markdown, **metadata},
                ensure_ascii=False,
            ),
            job["project_id"],
            job["stage"],
        ),
    )
    conn.execute(
        """
        update artifacts
        set content_json = %s::jsonb,
            updated_at = now()
        where project_id = %s and stage = %s and format = 'json'
        """,
        (
            json.dumps(
                {"stage": job["stage"], "markdown": markdown, **metadata},
                ensure_ascii=False,
            ),
            job["project_id"],
            job["stage"],
        ),
    )


def load_project(conn: psycopg.Connection, project_id: str) -> dict[str, Any]:
    row = conn.execute(
        "select * from projects where id = %s limit 1",
        (project_id,),
    ).fetchone()
    if row is None:
        raise RuntimeError(f"Project not found: {project_id}")
    return dict(row)


def load_artifact_markdown(conn: psycopg.Connection, project_id: str, stage: str) -> str:
    row = conn.execute(
        """
        select content_md
        from artifacts
        where project_id = %s and stage = %s and format = 'md'
        limit 1
        """,
        (project_id, stage),
    ).fetchone()
    if row is None:
        return ""
    content = row.get("content_md")
    if not isinstance(content, str):
        return ""
    return content.strip()


def load_source_files(conn: psycopg.Connection, project_id: str) -> list[dict[str, Any]]:
    rows = conn.execute(
        """
        select * from source_files
        where project_id = %s and upload_status = 'completed'
        order by position asc, created_at asc
        """,
        (project_id,),
    ).fetchall()
    return [dict(item) for item in rows]


def pick_video_source_file(files: list[dict[str, Any]]) -> dict[str, Any] | None:
    for source_file in files:
        mime_type = str(source_file.get("mime_type") or "").lower()
        kind = str(source_file.get("kind") or "").lower()
        if mime_type.startswith("video/") or kind == "video":
            return source_file
    return None


def file_extension(source_file: dict[str, Any]) -> str:
    name = str(source_file.get("original_name") or "").strip().lower()
    if "." not in name:
        return ""
    return name.rsplit(".", 1)[-1]


def is_document_source_file(source_file: dict[str, Any]) -> bool:
    mime_type = str(source_file.get("mime_type") or "").lower()
    kind = str(source_file.get("kind") or "").lower()
    ext = file_extension(source_file)
    if mime_type.startswith("video/") or kind == "video":
        return False
    if kind in {"document", "pdf", "presentation", "text"}:
        return True
    if mime_type.startswith("text/"):
        return True
    if mime_type in {
        "application/pdf",
        "application/rtf",
        "application/msword",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/vnd.ms-powerpoint",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }:
        return True
    return ext in {"pdf", "doc", "docx", "ppt", "pptx", "rtf", "txt", "md"}


def pick_document_source_files(files: list[dict[str, Any]]) -> list[dict[str, Any]]:
    return [source_file for source_file in files if is_document_source_file(source_file)]


def s3_client(config: WorkerConfig) -> Any:
    return boto3.client(
        "s3",
        endpoint_url=config.s3_endpoint,
        aws_access_key_id=config.s3_access_key_id,
        aws_secret_access_key=config.s3_secret_access_key,
        region_name=config.s3_region,
        config=Config(
            signature_version="s3v4",
            connect_timeout=10,
            read_timeout=60,
            retries={"max_attempts": 2, "mode": "standard"},
            s3={
                "addressing_style": "path",
                # Some S3-compatible providers reject signed payload hashing
                # for PUT requests; using unsigned payload avoids false mismatch.
                "payload_signing_enabled": False,
            },
        ),
    )


def s3_presigned_get_url(config: WorkerConfig, storage_key: str, expires_in_seconds: int) -> str:
    return s3_client(config).generate_presigned_url(
        "get_object",
        Params={"Bucket": config.s3_bucket, "Key": storage_key},
        ExpiresIn=max(60, expires_in_seconds),
    )


def log_worker_event(event: str, **fields: Any) -> None:
    payload: dict[str, Any] = {"service": "worker", "event": event}
    for key, value in fields.items():
        if value is not None:
            payload[key] = value
    print(json.dumps(payload, ensure_ascii=False))


def download_s3_object_bytes(config: WorkerConfig, storage_key: str) -> bytes:
    response = s3_client(config).get_object(Bucket=config.s3_bucket, Key=storage_key)
    body = response.get("Body")
    if body is None:
        raise RuntimeError("S3 вернул пустой body для source file.")
    content = body.read()
    if not isinstance(content, (bytes, bytearray)):
        raise RuntimeError("S3 вернул невалидный body для source file.")
    return bytes(content)


def download_s3_object_to_file(
    config: WorkerConfig,
    storage_key: str,
    destination_path: str,
    *,
    heartbeat_conn: psycopg.Connection,
    job_id: str,
    meeting_id: str,
) -> int:
    response = s3_client(config).get_object(Bucket=config.s3_bucket, Key=storage_key)
    body = response.get("Body")
    if body is None:
        raise RuntimeError("S3 вернул пустой body для source file.")

    bytes_written = 0
    chunk_size = 256 * 1024
    last_heartbeat_at = time.monotonic()

    with open(destination_path, "wb") as destination:
        while True:
            chunk = body.read(chunk_size)
            if not chunk:
                break
            destination.write(chunk)
            bytes_written += len(chunk)

            if time.monotonic() - last_heartbeat_at >= 10.0:
                touch_meeting_job_heartbeat(heartbeat_conn, job_id)
                log_worker_event(
                    "meeting-audio-download-heartbeat",
                    jobId=job_id,
                    meetingId=meeting_id,
                    storageKey=storage_key,
                    bytesSize=bytes_written,
                )
                last_heartbeat_at = time.monotonic()

    return bytes_written


def extract_document_text(source_file: dict[str, Any], file_bytes: bytes) -> str:
    mime_type = str(source_file.get("mime_type") or "").lower()
    ext = file_extension(source_file)
    if mime_type == "application/pdf" or ext == "pdf":
        reader = PdfReader(BytesIO(file_bytes))
        pages = [page.extract_text() or "" for page in reader.pages]
        return "\n\n".join(item.strip() for item in pages if item.strip())
    if (
        mime_type
        == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        or ext == "docx"
    ):
        document = Document(BytesIO(file_bytes))
        paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs]
        return "\n".join(item for item in paragraphs if item)
    if (
        mime_type
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        or ext == "pptx"
        or mime_type == "application/vnd.ms-powerpoint"
        or ext == "ppt"
    ):
        presentation = Presentation(BytesIO(file_bytes))
        slides_text: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            fragments: list[str] = []
            for shape in slide.shapes:
                text = str(getattr(shape, "text", "") or "").strip()
                if text:
                    fragments.append(text)
            if fragments:
                slides_text.append(f"Слайд {index}:\n" + "\n".join(fragments))
        return "\n\n".join(slides_text)
    if mime_type == "application/rtf" or ext == "rtf":
        decoded = file_bytes.decode("utf-8", errors="ignore")
        return rtf_to_text(decoded).strip()
    return file_bytes.decode("utf-8", errors="ignore").strip()


def call_transcription_service(
    config: WorkerConfig, source_file: dict[str, Any], timeout_seconds: float = 1200.0
) -> dict[str, Any]:
    endpoint = f"{config.transcription_service_url.rstrip('/')}/transcribe"
    payload = {
        "source": {
            "bucket": config.s3_bucket,
            "key": source_file["storage_key"],
            "originalName": source_file.get("original_name"),
            "mimeType": source_file.get("mime_type"),
        },
        "s3": {
            "endpoint": config.s3_endpoint,
            "bucket": config.s3_bucket,
            "region": config.s3_region,
            "accessKeyId": config.s3_access_key_id,
            "secretAccessKey": config.s3_secret_access_key,
        },
    }
    req = urlrequest.Request(
        endpoint,
        method="POST",
        data=json.dumps(payload, ensure_ascii=False).encode("utf-8"),
        headers={"Content-Type": "application/json; charset=utf-8"},
    )
    try:
        with urlrequest.urlopen(req, timeout=timeout_seconds) as response:
            raw = response.read().decode("utf-8")
    except urlerror.HTTPError as exc:
        detail = exc.read().decode("utf-8", errors="ignore")
        raise RuntimeError(
            f"Transcription service HTTP {exc.code}: {detail or exc.reason}"
        ) from exc
    except urlerror.URLError as exc:
        raise RuntimeError(f"Transcription service недоступен: {exc.reason}") from exc

    body = json.loads(raw)
    if not body.get("success"):
        error_payload = body.get("error") or {}
        message = error_payload.get("message") or "Ошибка транскрибации"
        raise RuntimeError(str(message))
    data = body.get("data")
    if not isinstance(data, dict):
        raise RuntimeError("Transcription service вернул невалидный payload")
    return data


def load_job(conn: psycopg.Connection, job_id: str) -> dict[str, Any] | None:
    row = conn.execute(
        "select * from processing_jobs where id = %s limit 1",
        (job_id,),
    ).fetchone()
    return dict(row) if row else None


def recover_missing_job(
    conn: psycopg.Connection, job_message: dict[str, Any]
) -> dict[str, Any] | None:
    job_id = str(job_message.get("jobId") or "").strip()
    project_id = str(job_message.get("projectId") or "").strip()
    stage = str(job_message.get("stage") or "").strip()
    trigger = str(job_message.get("trigger") or "auto").strip() or "auto"

    if not job_id or not project_id or not stage:
        return None

    if stage not in STAGE_ORDER:
        return None

    prompt_keys = [prompt.key for prompt in prompt_bundle_for_stage(stage)]

    conn.execute(
        """
        insert into processing_jobs (
          id, project_id, stage, status, payload_json, result_json, error_text, started_at, finished_at, created_at
        ) values (
          %s, %s, %s, 'queued', %s::jsonb, null, null, null, null, now()
        )
        on conflict (id) do nothing
        """,
        (
            job_id,
            project_id,
            stage,
            json.dumps(
                {
                    "stage": stage,
                    "trigger": trigger,
                    "promptKeys": prompt_keys,
                    "autoGenerateAll": trigger == "auto",
                    "nextStage": None,
                },
                ensure_ascii=False,
            ),
        ),
    )
    return load_job(conn, job_id)


def append_project_log(conn: psycopg.Connection, project_id: str, message: str) -> None:
    row = conn.execute(
        "select logs from projects where id = %s limit 1",
        (project_id,),
    ).fetchone()
    logs = []
    if row and row["logs"]:
        logs = list(row["logs"])
    logs = [message, *logs][:10]
    conn.execute(
        "update projects set logs = %s::jsonb, updated_at = now() where id = %s",
        (json.dumps(logs, ensure_ascii=False), project_id),
    )


def set_job_processing(conn: psycopg.Connection, job_id: str) -> None:
    conn.execute(
        """
        update processing_jobs
        set status = 'processing', started_at = now(), error_text = null
        where id = %s
        """,
        (job_id,),
    )


def set_job_done(
    conn: psycopg.Connection, job_id: str, stage: str, metadata: dict[str, Any]
) -> None:
    conn.execute(
        """
        update processing_jobs
        set status = 'done',
            result_json = %s::jsonb,
            finished_at = now()
        where id = %s
        """,
        (
            json.dumps(
                {
                    "status": "done",
                    "stage": stage,
                    "generatedAt": utc_now(),
                    **metadata,
                },
                ensure_ascii=False,
            ),
            job_id,
        ),
    )


def set_job_failed(conn: psycopg.Connection, job_id: str, error_text: str) -> None:
    conn.execute(
        """
        update processing_jobs
        set status = 'failed',
            error_text = %s,
            finished_at = now()
        where id = %s
        """,
        (error_text, job_id),
    )


def salutespeech_request(
    url: str,
    *,
    method: str = "GET",
    headers: dict[str, str] | None = None,
    data: bytes | None = None,
    timeout: float = 60.0,
    verify_ssl: bool = True,
    ca_cert_path: str = "",
) -> Any:
    request = urlrequest.Request(
        url,
        method=method,
        data=data,
        headers=headers or {},
    )
    if not verify_ssl:
        ssl_context = ssl._create_unverified_context()
    elif ca_cert_path:
        ssl_context = ssl.create_default_context(cafile=ca_cert_path)
    else:
        ssl_context = ssl.create_default_context()
    try:
        with urlrequest.urlopen(request, timeout=timeout, context=ssl_context) as response:
            raw_bytes = response.read()
            content_type = str(response.headers.get("Content-Type") or "")
            raw = raw_bytes.decode("utf-8", errors="replace")
    except urlerror.HTTPError as exc:
        try:
            details = exc.read().decode("utf-8", errors="ignore").strip()
        except Exception:
            details = ""
        suffix = f": {details}" if details else ""
        message = normalize_provider_error(
            "salutespeech",
            f"HTTP {exc.code} from {url}{suffix}",
        )
        raise RuntimeError(message) from exc
    try:
        payload = json.loads(raw)
    except json.JSONDecodeError as exc:
        snippet = raw[:500].replace("\n", "\\n")
        suffix = f" content-type={content_type}" if content_type else ""
        raise RuntimeError(
            f"SaluteSpeech вернул невалидный JSON.{suffix} body={snippet}"
        ) from exc
    if not isinstance(payload, (dict, list)):
        raise RuntimeError("SaluteSpeech вернул невалидный JSON-тип.")
    return payload


def get_salutespeech_access_token(config: WorkerConfig) -> str:
    if not config.salutespeech_auth_key:
        raise RuntimeError("SALUTESPEECH_AUTH_KEY не задан.")
    if not config.salutespeech_oauth_url:
        raise RuntimeError("SALUTESPEECH_OAUTH_URL не задан.")

    data = f"scope={config.salutespeech_scope}".encode("utf-8")
    payload = salutespeech_request(
        config.salutespeech_oauth_url,
        method="POST",
        data=data,
        headers={
            "Authorization": f"Basic {config.salutespeech_auth_key}",
            "RqUID": str(uuid.uuid4()),
            "Content-Type": "application/x-www-form-urlencoded",
        },
        timeout=30.0,
        verify_ssl=not config.salutespeech_ssl_no_verify,
        ca_cert_path=config.salutespeech_ca_cert_path,
    )
    if not isinstance(payload, dict):
        raise RuntimeError("SaluteSpeech OAuth вернул неожиданный формат ответа.")
    token = str(
        payload.get("access_token")
        or payload.get("result", {}).get("access_token")
        or ""
    ).strip()
    if not token:
        raise RuntimeError("SaluteSpeech не вернул access token.")
    return token


def encode_multipart_form(field_name: str, file_name: str, content: bytes) -> tuple[bytes, str]:
    boundary = f"----ecolms-boundary-{int(time.time() * 1000)}"
    body = (
        f"--{boundary}\r\n"
        f'Content-Disposition: form-data; name="{field_name}"; filename="{file_name}"\r\n'
        "Content-Type: application/octet-stream\r\n\r\n"
    ).encode("utf-8") + content + f"\r\n--{boundary}--\r\n".encode("utf-8")
    return body, boundary


def salutespeech_upload_file(
    config: WorkerConfig, token: str, file_name: str, content: bytes
) -> str:
    body, boundary = encode_multipart_form("file", file_name, content)
    payload = salutespeech_request(
        config.salutespeech_upload_url,
        method="POST",
        data=body,
        headers={
            "Authorization": f"Bearer {token}",
            "Content-Type": f"multipart/form-data; boundary={boundary}",
        },
        timeout=120.0,
        verify_ssl=not config.salutespeech_ssl_no_verify,
        ca_cert_path=config.salutespeech_ca_cert_path,
    )
    if not isinstance(payload, dict):
        raise RuntimeError("SaluteSpeech upload вернул неожиданный формат ответа.")
    request_file_id = str(
        payload.get("result", {}).get("request_file_id")
        or payload.get("request_file_id")
        or ""
    ).strip()
    if not request_file_id:
        raise RuntimeError("SaluteSpeech не вернул request_file_id.")
    return request_file_id


def salutespeech_create_recognition_task(
    config: WorkerConfig,
    token: str,
    *,
    request_file_id: str,
    audio_encoding: str,
    sample_rate: int,
    channels_count: int,
) -> str:
    options: dict[str, Any] = {
        "model": config.salutespeech_model,
        "language": config.salutespeech_language,
        "audio_encoding": audio_encoding,
        "channels_count": channels_count,
        "speaker_separation_options": {
            "enable": True,
        },
    }
    if audio_encoding in {"PCM_S16LE", "ALAW", "MULAW"}:
        options["sample_rate"] = sample_rate

    payload = salutespeech_request(
        config.salutespeech_recognize_url,
        method="POST",
        data=json.dumps(
            {
                "request_file_id": request_file_id,
                "options": options,
            },
            ensure_ascii=False,
        ).encode("utf-8"),
        headers={
            "Authorization": f"Bearer {token}",
            "Content-Type": "application/json; charset=utf-8",
        },
        timeout=60.0,
        verify_ssl=not config.salutespeech_ssl_no_verify,
        ca_cert_path=config.salutespeech_ca_cert_path,
    )
    if not isinstance(payload, dict):
        raise RuntimeError("SaluteSpeech recognize вернул неожиданный формат ответа.")
    task_id = str(payload.get("result", {}).get("id") or payload.get("id") or "").strip()
    if not task_id:
        raise RuntimeError("SaluteSpeech не вернул id задачи распознавания.")
    return task_id


def salutespeech_get_task_status(config: WorkerConfig, token: str, task_id: str) -> dict[str, Any]:
    payload = salutespeech_request(
        f"{config.salutespeech_task_url}?id={task_id}",
        headers={"Authorization": f"Bearer {token}"},
        timeout=30.0,
        verify_ssl=not config.salutespeech_ssl_no_verify,
        ca_cert_path=config.salutespeech_ca_cert_path,
    )
    if not isinstance(payload, dict):
        raise RuntimeError("SaluteSpeech task:get вернул неожиданный формат ответа.")
    return payload


def salutespeech_download_result(
    config: WorkerConfig, token: str, response_file_id: str
) -> Any:
    return salutespeech_request(
        f"{config.salutespeech_download_url}?response_file_id={response_file_id}",
        headers={"Authorization": f"Bearer {token}"},
        timeout=120.0,
        verify_ssl=not config.salutespeech_ssl_no_verify,
        ca_cert_path=config.salutespeech_ca_cert_path,
    )


def extract_float(value: Any) -> float | None:
    if value is None:
        return None
    if isinstance(value, (int, float)):
        return float(value)
    if isinstance(value, str):
        normalized = value.strip().lower()
        if normalized.endswith("ms"):
            try:
                return float(normalized[:-2].strip()) / 1000.0
            except ValueError:
                return None
        if normalized.endswith("s"):
            try:
                return float(normalized[:-1].strip())
            except ValueError:
                return None
        try:
            return float(normalized)
        except ValueError:
            return None
    return None


def to_milliseconds(value: Any) -> int | None:
    number = extract_float(value)
    if number is None:
        return None
    if number > 10_000:
        return int(number)
    return int(number * 1000)


def collect_segment_candidates(value: Any, results: list[dict[str, Any]]) -> None:
    if isinstance(value, dict):
        has_speaker = "speaker_id" in value
        has_text = isinstance(value.get("text"), str) and value.get("text", "").strip()
        has_word = isinstance(value.get("word"), str) and value.get("word", "").strip()
        if has_speaker and (has_text or has_word):
            results.append(value)
        for nested in value.values():
            collect_segment_candidates(nested, results)
    elif isinstance(value, list):
        for item in value:
            collect_segment_candidates(item, results)


def collect_chunk_style_candidates(value: Any, results: list[dict[str, Any]]) -> None:
    if isinstance(value, dict):
        speaker_info = value.get("speaker_info")
        speaker_id: int | None = None
        if isinstance(speaker_info, dict):
            speaker_raw = speaker_info.get("speaker_id")
            if isinstance(speaker_raw, int):
                speaker_id = speaker_raw
            elif isinstance(speaker_raw, str):
                try:
                    speaker_id = int(speaker_raw)
                except ValueError:
                    speaker_id = None

        hypotheses = value.get("results")
        if isinstance(hypotheses, list) and speaker_id is not None and speaker_id >= 0:
            for item in hypotheses:
                if not isinstance(item, dict):
                    continue
                text = str(item.get("text") or "").strip()
                if text:
                    results.append(
                        {
                            "speaker_id": speaker_id,
                            "text": text,
                            "start": item.get("start"),
                            "end": item.get("end"),
                            "confidence": item.get("confidence"),
                        }
                    )
                word_alignments = item.get("word_alignments")
                if isinstance(word_alignments, list):
                    for word_item in word_alignments:
                        if not isinstance(word_item, dict):
                            continue
                        word = str(word_item.get("word") or "").strip()
                        if not word:
                            continue
                        results.append(
                            {
                                "speaker_id": speaker_id,
                                "word": word,
                                "start": word_item.get("start"),
                                "end": word_item.get("end"),
                                "confidence": word_item.get("confidence"),
                            }
                        )

        for nested in value.values():
            collect_chunk_style_candidates(nested, results)
    elif isinstance(value, list):
        for item in value:
            collect_chunk_style_candidates(item, results)


def parse_salutespeech_segments(payload: Any) -> list[dict[str, Any]]:
    candidates: list[dict[str, Any]] = []
    collect_segment_candidates(payload, candidates)
    collect_chunk_style_candidates(payload, candidates)
    segment_like: list[dict[str, Any]] = []
    word_like: list[dict[str, Any]] = []

    for item in candidates:
        speaker_id = int(item.get("speaker_id", -1))
        if speaker_id < 0:
            continue

        if isinstance(item.get("text"), str) and item.get("text", "").strip():
            start_ms = (
                to_milliseconds(item.get("start_ms"))
                or to_milliseconds(item.get("start"))
                or to_milliseconds(item.get("start_time"))
            )
            end_ms = (
                to_milliseconds(item.get("end_ms"))
                or to_milliseconds(item.get("end"))
                or to_milliseconds(item.get("end_time"))
            )
            if start_ms is None or end_ms is None or end_ms < start_ms:
                continue
            segment_like.append(
                {
                    "speaker_id": speaker_id,
                    "text": str(item.get("text")).strip(),
                    "start_ms": start_ms,
                    "end_ms": end_ms,
                    "confidence": extract_float(item.get("confidence")),
                    "provider_payload_json": item,
                }
            )
        elif isinstance(item.get("word"), str) and item.get("word", "").strip():
            start_ms = (
                to_milliseconds(item.get("start_ms"))
                or to_milliseconds(item.get("start"))
                or to_milliseconds(item.get("start_time"))
            )
            end_ms = (
                to_milliseconds(item.get("end_ms"))
                or to_milliseconds(item.get("end"))
                or to_milliseconds(item.get("end_time"))
            )
            if start_ms is None or end_ms is None or end_ms < start_ms:
                continue
            word_like.append(
                {
                    "speaker_id": speaker_id,
                    "word": str(item.get("word")).strip(),
                    "start_ms": start_ms,
                    "end_ms": end_ms,
                    "confidence": extract_float(item.get("confidence")),
                    "provider_payload_json": item,
                }
            )

    if segment_like:
        segment_like.sort(key=lambda item: (item["start_ms"], item["end_ms"]))
        return segment_like

    if not word_like:
        raise RuntimeError("SaluteSpeech не вернул diarized segments со speaker_id.")

    word_like.sort(key=lambda item: (item["start_ms"], item["end_ms"]))
    merged: list[dict[str, Any]] = []
    current: dict[str, Any] | None = None
    for word in word_like:
        if current is None:
            current = {
                "speaker_id": word["speaker_id"],
                "text": word["word"],
                "start_ms": word["start_ms"],
                "end_ms": word["end_ms"],
                "confidence_values": [word.get("confidence")],
                "provider_items": [word["provider_payload_json"]],
            }
            continue
        gap = int(word["start_ms"]) - int(current["end_ms"])
        if (
            int(word["speaker_id"]) == int(current["speaker_id"])
            and gap <= 1200
        ):
            current["text"] = f"{current['text']} {word['word']}".strip()
            current["end_ms"] = word["end_ms"]
            current["confidence_values"].append(word.get("confidence"))
            current["provider_items"].append(word["provider_payload_json"])
            continue
        confidences = [value for value in current["confidence_values"] if value is not None]
        merged.append(
            {
                "speaker_id": current["speaker_id"],
                "text": current["text"].strip(),
                "start_ms": current["start_ms"],
                "end_ms": current["end_ms"],
                "confidence": round(sum(confidences) / len(confidences), 4)
                if confidences
                else None,
                "provider_payload_json": {"words": current["provider_items"]},
            }
        )
        current = {
            "speaker_id": word["speaker_id"],
            "text": word["word"],
            "start_ms": word["start_ms"],
            "end_ms": word["end_ms"],
            "confidence_values": [word.get("confidence")],
            "provider_items": [word["provider_payload_json"]],
        }

    if current is not None:
        confidences = [value for value in current["confidence_values"] if value is not None]
        merged.append(
            {
                "speaker_id": current["speaker_id"],
                "text": current["text"].strip(),
                "start_ms": current["start_ms"],
                "end_ms": current["end_ms"],
                "confidence": round(sum(confidences) / len(confidences), 4)
                if confidences
                else None,
                "provider_payload_json": {"words": current["provider_items"]},
            }
        )

    return merged


def format_ms_range(start_ms: int, end_ms: int) -> str:
    def one(value: int) -> str:
        seconds_total = max(0, value // 1000)
        hours = seconds_total // 3600
        minutes = (seconds_total % 3600) // 60
        seconds = seconds_total % 60
        return f"{hours:02d}:{minutes:02d}:{seconds:02d}"

    return f"{one(start_ms)} - {one(end_ms)}"


def build_meeting_transcript_markdown(title: str, segments: list[dict[str, Any]]) -> str:
    lines = [f"# {title}", ""]
    for segment in segments:
        lines.append(
            f"## [{format_ms_range(int(segment['start_ms']), int(segment['end_ms']))}] {segment['display_name']}"
        )
        lines.append(str(segment["text"]).strip())
        lines.append("")
    return "\n".join(lines).strip() + "\n"


def build_meeting_transcript_json(
    meeting_id: str,
    speakers: list[dict[str, Any]],
    segments: list[dict[str, Any]],
) -> dict[str, Any]:
    def normalize_confidence(value: Any) -> float | None:
        if value is None:
            return None
        if isinstance(value, (int, float)):
            return float(value)
        try:
            return float(value)
        except (TypeError, ValueError):
            return None

    return {
        "meetingId": meeting_id,
        "speakers": [
            {
                "id": item["id"],
                "speakerLabel": item["speaker_label"],
                "displayName": item["display_name"],
                "isUserEdited": bool(item["is_user_edited"]),
                "sortOrder": int(item["sort_order"]),
            }
            for item in speakers
        ],
        "segments": [
            {
                "id": int(item["id"]),
                "speakerId": item["speaker_id"],
                "speakerLabel": item["speaker_label"],
                "displayName": item["display_name"],
                "startMs": int(item["start_ms"]),
                "endMs": int(item["end_ms"]),
                "text": str(item["text"]),
                "confidence": normalize_confidence(item["confidence"]),
            }
            for item in segments
        ],
    }


def build_meeting_llm_transcript_input(segments: list[dict[str, Any]]) -> str:
    parts: list[str] = []
    for segment in segments:
        parts.append(
            "\n".join(
                [
                    f"[segment_id={segment['id']}]",
                    f"[{format_ms_range(int(segment['start_ms']), int(segment['end_ms']))}] {segment['display_name']}",
                    str(segment["text"]).strip(),
                ]
            )
        )
    return "\n\n".join(parts).strip()


def build_meeting_actions_markdown(result: dict[str, Any]) -> str:
    def extract_text(item: Any) -> str:
        if isinstance(item, str):
            return item.strip()
        if isinstance(item, dict):
            value = str(item.get("text") or item.get("title") or "").strip()
            return value
        return ""

    def extract_list(value: Any) -> list[dict[str, Any] | str]:
        if not isinstance(value, list):
            return []
        normalized: list[dict[str, Any] | str] = []
        for item in value:
            if isinstance(item, (dict, str)):
                normalized.append(item)
        return normalized

    decisions = extract_list(result.get("decisions"))
    action_items = extract_list(result.get("actionItems"))
    open_questions = extract_list(result.get("openQuestions"))
    short_summary = str(result.get("shortSummary") or "").strip()

    lines: list[str] = []
    if short_summary:
        lines.extend(["# Поручения и решения", "", short_summary, ""])

    if decisions:
        lines.extend(["## Решения", ""])
        for item in decisions:
            text = extract_text(item)
            if text:
                lines.append(f"- {text}")
        lines.append("")

    if action_items:
        lines.extend(["## Поручения", ""])
        for item in action_items:
            if not isinstance(item, dict):
                text = extract_text(item)
                if text:
                    lines.append(f"- {text}")
                continue
            text = extract_text(item)
            if not text:
                continue
            details: list[str] = []
            assignee = str(item.get("assignee") or "").strip()
            deadline = str(item.get("deadline") or "").strip()
            source_segment_ids = item.get("sourceSegmentIds")
            if assignee:
                details.append(f"ответственный: {assignee}")
            if deadline:
                details.append(f"дедлайн: {deadline}")
            if isinstance(source_segment_ids, list):
                normalized_ids = [str(value).strip() for value in source_segment_ids if str(value).strip()]
                if normalized_ids:
                    details.append(f"segment_id: {', '.join(normalized_ids)}")
            if details:
                lines.append(f"- {text} ({'; '.join(details)})")
            else:
                lines.append(f"- {text}")
        lines.append("")

    if open_questions:
        lines.extend(["## Открытые вопросы", ""])
        for item in open_questions:
            text = extract_text(item)
            if text:
                lines.append(f"- {text}")
        lines.append("")

    markdown = "\n".join(lines).strip()
    return f"{markdown}\n" if markdown else ""


def generate_meeting_markdown_with_llm(
    conn: psycopg.Connection,
    config: WorkerConfig,
    *,
    stage: str,
    meeting_title: str,
    transcript_input: str,
) -> dict[str, Any]:
    if not transcript_input.strip():
        raise RuntimeError(f"Невозможно вызвать LLM для {stage}: пустой transcript.")
    prompt_definition = load_prompt_definition(conn, "meetings", stage)

    user_payload = {
        "meetingTitle": meeting_title,
        "stage": stage,
        "task": prompt_definition.user_prompt_template,
        "transcript": transcript_input[:180000],
    }
    messages = [
        {
            "role": "system",
            "content": prompt_definition.system_prompt,
        },
        {"role": "user", "content": json.dumps(user_payload, ensure_ascii=False)},
    ]

    provider = resolve_llm_provider(config)

    try:
        if provider == "openai":
            content, model_name = openai_chat_completion(
                config.openai_api_key,
                config.openai_model,
                messages,
                config.llm_timeout_seconds,
            )
        else:
            content, model_name = openrouter_chat_completion(
                config.openrouter_api_key,
                config.openrouter_base_url,
                config.openrouter_model,
                messages,
                config.llm_timeout_seconds,
            )
        parsed = parse_json_from_text(content)
        markdown = str(parsed.get("markdown") or "").strip()
        short_summary = str(parsed.get("shortSummary") or "").strip()
        if stage == "meeting_actions" and not markdown:
            markdown = build_meeting_actions_markdown(parsed).strip()
        if not markdown:
            raise RuntimeError("LLM вернул пустое поле markdown.")
        return {
            **parsed,
            "markdown": markdown,
            "shortSummary": short_summary,
            "provider": provider,
            "model": model_name,
            "attempts": [{"provider": provider, "ok": True}],
            "promptModule": prompt_definition.module,
            "promptKey": prompt_definition.key,
            "promptTitle": prompt_definition.title,
        }
    except Exception as exc:
        raise RuntimeError(f"LLM-провайдер {provider} завершился ошибкой: {exc}") from exc


def meeting_actions_json_from_llm(result: dict[str, Any], markdown: str) -> dict[str, Any]:
    decisions = result.get("decisions")
    action_items = result.get("actionItems")
    open_questions = result.get("openQuestions")
    return {
        "decisions": decisions if isinstance(decisions, list) else [],
        "actionItems": action_items if isinstance(action_items, list) else [],
        "openQuestions": open_questions if isinstance(open_questions, list) else [],
        "markdown": markdown,
    }


def load_meeting(conn: psycopg.Connection, meeting_id: str) -> dict[str, Any]:
    row = conn.execute(
        "select * from meetings where id = %s limit 1",
        (meeting_id,),
    ).fetchone()
    if row is None:
        raise RuntimeError(f"Meeting not found: {meeting_id}")
    return dict(row)


def load_meeting_job(conn: psycopg.Connection, job_id: str) -> dict[str, Any] | None:
    row = conn.execute(
        "select * from meeting_jobs where id = %s limit 1",
        (job_id,),
    ).fetchone()
    return dict(row) if row else None


def load_meeting_source_file(conn: psycopg.Connection, meeting_id: str) -> dict[str, Any] | None:
    row = conn.execute(
        """
        select *
        from meeting_source_files
        where meeting_id = %s
        limit 1
        """,
        (meeting_id,),
    ).fetchone()
    return dict(row) if row else None


def set_meeting_job_processing(conn: psycopg.Connection, job_id: str) -> None:
    conn.execute(
        """
        update meeting_jobs
        set
          status = 'processing',
          started_at = now(),
          processing_heartbeat_at = now(),
          error_text = null
        where id = %s
        """,
        (job_id,),
    )


def touch_meeting_job_heartbeat(conn: psycopg.Connection, job_id: str) -> None:
    conn.execute(
        """
        update meeting_jobs
        set processing_heartbeat_at = now()
        where id = %s and status = 'processing'
        """,
        (job_id,),
    )


def set_meeting_job_done(
    conn: psycopg.Connection, job_id: str, stage: str, metadata: dict[str, Any]
) -> None:
    conn.execute(
        """
        update meeting_jobs
        set status = 'done',
            result_json = %s::jsonb,
            processing_heartbeat_at = now(),
            finished_at = now()
        where id = %s
        """,
        (
            json.dumps(
                {"status": "done", "stage": stage, "generatedAt": utc_now(), **metadata},
                ensure_ascii=False,
            ),
            job_id,
        ),
    )


def set_meeting_job_failed(conn: psycopg.Connection, job_id: str, error_text: str) -> None:
    conn.execute(
        """
        update meeting_jobs
        set status = 'failed',
            error_text = %s,
            processing_heartbeat_at = now(),
            finished_at = now()
        where id = %s
        """,
        (error_text, job_id),
    )


def ensure_meeting_artifact(
    conn: psycopg.Connection,
    meeting_id: str,
    stage: str,
    markdown: str,
    content_json: dict[str, Any],
) -> None:
    conn.execute(
        """
        update meeting_artifacts
        set content_md = %s,
            updated_at = now()
        where meeting_id = %s and stage = %s and format = 'md'
        """,
        (markdown, meeting_id, stage),
    )
    conn.execute(
        """
        update meeting_artifacts
        set content_json = %s::jsonb,
            updated_at = now()
        where meeting_id = %s and stage = %s and format = 'json'
        """,
        (json.dumps(content_json, ensure_ascii=False), meeting_id, stage),
    )


def load_meeting_artifact(
    conn: psycopg.Connection, meeting_id: str, stage: str, format_name: str
) -> dict[str, Any] | None:
    row = conn.execute(
        """
        select *
        from meeting_artifacts
        where meeting_id = %s and stage = %s and format = %s
        limit 1
        """,
        (meeting_id, stage, format_name),
    ).fetchone()
    return dict(row) if row else None


def load_meeting_speakers(conn: psycopg.Connection, meeting_id: str) -> list[dict[str, Any]]:
    rows = conn.execute(
        """
        select *
        from meeting_speakers
        where meeting_id = %s
        order by sort_order asc, created_at asc
        """,
        (meeting_id,),
    ).fetchall()
    return [dict(item) for item in rows]


def load_meeting_segments(conn: psycopg.Connection, meeting_id: str) -> list[dict[str, Any]]:
    rows = conn.execute(
        """
        select *
        from meeting_speaker_segments
        where meeting_id = %s
        order by start_ms asc, id asc
        """,
        (meeting_id,),
    ).fetchall()
    return [dict(item) for item in rows]


def clear_meeting_transcript(conn: psycopg.Connection, meeting_id: str) -> None:
    conn.execute("delete from meeting_speaker_segments where meeting_id = %s", (meeting_id,))
    conn.execute("delete from meeting_speakers where meeting_id = %s", (meeting_id,))


def upsert_meeting_transcript(
    conn: psycopg.Connection, meeting_id: str, raw_segments: list[dict[str, Any]]
) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
    clear_meeting_transcript(conn, meeting_id)
    speaker_map: dict[int, dict[str, Any]] = {}
    sort_order = 0
    for segment in raw_segments:
        speaker_key = int(segment["speaker_id"])
        if speaker_key in speaker_map:
            continue
        sort_order += 1
        speaker_id = f"{meeting_id}-speaker-{speaker_key}"
        speaker = {
            "id": speaker_id,
            "meeting_id": meeting_id,
            "speaker_label": f"speaker_{speaker_key}",
            "display_name": f"Спикер {sort_order}",
            "is_user_edited": False,
            "sort_order": sort_order,
        }
        speaker_map[speaker_key] = speaker
        conn.execute(
            """
            insert into meeting_speakers (
              id, meeting_id, speaker_label, display_name, is_user_edited, sort_order, created_at, updated_at
            ) values (
              %s, %s, %s, %s, %s, %s, now(), now()
            )
            """,
            (
                speaker["id"],
                speaker["meeting_id"],
                speaker["speaker_label"],
                speaker["display_name"],
                speaker["is_user_edited"],
                speaker["sort_order"],
            ),
        )

    inserted_segments: list[dict[str, Any]] = []
    for segment in raw_segments:
        speaker = speaker_map[int(segment["speaker_id"])]
        row = conn.execute(
            """
            insert into meeting_speaker_segments (
              meeting_id, speaker_id, speaker_label, start_ms, end_ms, text, confidence, provider_payload_json, created_at
            ) values (
              %s, %s, %s, %s, %s, %s, %s, %s::jsonb, now()
            )
            returning *
            """,
            (
                meeting_id,
                speaker["id"],
                speaker["speaker_label"],
                int(segment["start_ms"]),
                int(segment["end_ms"]),
                str(segment["text"]).strip(),
                segment.get("confidence"),
                json.dumps(segment.get("provider_payload_json") or {}, ensure_ascii=False),
            ),
        ).fetchone()
        if row is not None:
            inserted = dict(row)
            inserted["display_name"] = speaker["display_name"]
            inserted_segments.append(inserted)

    speakers = list(speaker_map.values())
    return speakers, inserted_segments


def ffprobe_audio_details(source_path: str) -> tuple[str, int, float]:
    try:
        result = subprocess.run(
            [
                "ffprobe",
                "-v",
                "error",
                "-select_streams",
                "a:0",
                "-show_entries",
                "stream=codec_name,channels",
                "-show_entries",
                "format=duration",
                "-of",
                "json",
                source_path,
            ],
            capture_output=True,
            text=True,
            timeout=60.0,
        )
    except subprocess.TimeoutExpired as exc:
        raise RuntimeError("ffprobe timed out while preparing meeting audio.") from exc
    if result.returncode != 0:
        raise RuntimeError(
            f"ffprobe failed: {(result.stderr or result.stdout or '').strip()}"
        )
    payload = json.loads(result.stdout or "{}")
    stream = payload.get("streams", [{}])[0]
    codec_name = str(stream.get("codec_name") or "").strip().lower()
    channels = int(stream.get("channels") or 1)
    duration = float(payload.get("format", {}).get("duration") or 0.0)
    return codec_name, channels, duration


def run_ffmpeg_with_heartbeat(
    command: list[str],
    *,
    timeout_seconds: float,
    heartbeat_conn: psycopg.Connection,
    job_id: str,
    meeting_id: str,
    storage_key: str | None,
) -> subprocess.CompletedProcess[str]:
    process = subprocess.Popen(
        command,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
    )
    heartbeat_interval_seconds = 10.0
    deadline = time.monotonic() + timeout_seconds

    while True:
        remaining = deadline - time.monotonic()
        if remaining <= 0:
            process.kill()
            stdout, stderr = process.communicate()
            raise RuntimeError("FFmpeg audio preparation timed out.")

        try:
            stdout, stderr = process.communicate(
                timeout=min(heartbeat_interval_seconds, remaining)
            )
            return subprocess.CompletedProcess(command, process.returncode, stdout, stderr)
        except subprocess.TimeoutExpired:
            touch_meeting_job_heartbeat(heartbeat_conn, job_id)
            log_worker_event(
                "meeting-audio-ffmpeg-heartbeat",
                jobId=job_id,
                meetingId=meeting_id,
                storageKey=storage_key,
            )


def upload_s3_object_bytes(
    config: WorkerConfig, storage_key: str, content: bytes, content_type: str
) -> None:
    s3_client(config).put_object(
        Bucket=config.s3_bucket,
        Key=storage_key,
        Body=content,
        ContentType=content_type,
    )


def prepare_meeting_audio(
    config: WorkerConfig,
    source_file: dict[str, Any],
    heartbeat_conn: psycopg.Connection,
    job_id: str,
) -> dict[str, Any]:
    log_worker_event(
        "meeting-audio-download-start",
        meetingId=source_file.get("meeting_id"),
        storageKey=source_file.get("storage_key"),
    )
    source_suffix = "." + file_extension(source_file) if file_extension(source_file) else ".bin"
    with tempfile.NamedTemporaryFile(suffix=source_suffix, delete=False) as src:
        source_path = src.name
    source_bytes_size = download_s3_object_to_file(
        config,
        str(source_file["storage_key"]),
        source_path,
        heartbeat_conn=heartbeat_conn,
        job_id=job_id,
        meeting_id=str(source_file.get("meeting_id")),
    )
    log_worker_event(
        "meeting-audio-download-done",
        meetingId=source_file.get("meeting_id"),
        storageKey=source_file.get("storage_key"),
        bytesSize=source_bytes_size,
    )
    with tempfile.NamedTemporaryFile(suffix=".ogg", delete=False) as out:
        output_path = out.name

    try:
        log_worker_event(
            "meeting-audio-ffprobe-start",
            meetingId=source_file.get("meeting_id"),
            storageKey=source_file.get("storage_key"),
        )
        source_codec, target_channels, duration = ffprobe_audio_details(source_path)
        log_worker_event(
            "meeting-audio-ffprobe-done",
            meetingId=source_file.get("meeting_id"),
            storageKey=source_file.get("storage_key"),
            sourceCodec=source_codec,
            durationSeconds=duration,
            targetChannels=target_channels,
        )
        fast_path_copy = source_codec == "opus" and target_channels <= 2
        if fast_path_copy:
            command = [
                "ffmpeg",
                "-hide_banner",
                "-loglevel",
                "error",
                "-y",
                "-i",
                source_path,
                "-vn",
                "-c:a",
                "copy",
                "-f",
                "ogg",
                output_path,
            ]
        else:
            command = [
                "ffmpeg",
                "-hide_banner",
                "-loglevel",
                "error",
                "-y",
                "-i",
                source_path,
                "-vn",
                "-ac",
                str(target_channels),
                "-ar",
                "48000",
                "-c:a",
                "libopus",
                "-b:a",
                "48k",
                "-f",
                "ogg",
                output_path,
            ]
        log_worker_event(
            "meeting-audio-ffmpeg-start",
            meetingId=source_file.get("meeting_id"),
            storageKey=source_file.get("storage_key"),
            timeoutSeconds=config.meeting_audio_prep_timeout_seconds,
            mode="copy" if fast_path_copy else "transcode",
        )
        result = run_ffmpeg_with_heartbeat(
            command,
            timeout_seconds=config.meeting_audio_prep_timeout_seconds,
            heartbeat_conn=heartbeat_conn,
            job_id=job_id,
            meeting_id=str(source_file.get("meeting_id")),
            storage_key=str(source_file.get("storage_key")) if source_file.get("storage_key") else None,
        )
        if result.returncode != 0:
            details = (result.stderr or result.stdout or "").strip()
            raise RuntimeError(f"FFmpeg audio preparation failed: {details or 'unknown error'}")
        log_worker_event(
            "meeting-audio-ffmpeg-done",
            meetingId=source_file.get("meeting_id"),
            storageKey=source_file.get("storage_key"),
            mode="copy" if fast_path_copy else "transcode",
        )
        with open(output_path, "rb") as file:
            audio_bytes = file.read()
        audio_storage_key: str | None = f"meetings/{source_file['meeting_id']}/derived/prepared-audio.ogg"
        try:
            upload_s3_object_bytes(
                config,
                audio_storage_key,
                audio_bytes,
                "audio/ogg",
            )
        except Exception as upload_error:
            # Prepared audio is an auxiliary artifact. If provider-specific
            # S3 upload fails, keep processing with in-memory bytes.
            print(
                json.dumps(
                    {
                        "service": "worker",
                        "status": "warning",
                        "scope": "meetings",
                        "event": "prepared-audio-upload-failed",
                        "meetingId": source_file.get("meeting_id"),
                        "error": str(upload_error),
                    },
                    ensure_ascii=False,
                )
            )
            audio_storage_key = None
        log_worker_event(
            "meeting-audio-ready",
            meetingId=source_file.get("meeting_id"),
            storageKey=audio_storage_key,
            durationSeconds=int(round(duration)) if duration > 0 else None,
            channelsCount=target_channels,
        )
        return {
            "audio_storage_key": audio_storage_key,
            "audio_mime_type": "audio/ogg",
            "duration_seconds": int(round(duration)) if duration > 0 else None,
            "sample_rate": 48000,
            "channels_count": target_channels,
            "audio_encoding": "OPUS",
            "audio_bytes": audio_bytes,
            "file_name": "prepared-audio.ogg",
        }
    finally:
        try:
            os.unlink(source_path)
        except OSError:
            pass
        try:
            os.unlink(output_path)
        except OSError:
            pass


def assemblyai_request(
    url: str,
    *,
    method: str = "GET",
    headers: dict[str, str] | None = None,
    data: bytes | None = None,
    timeout: float = 60.0,
) -> dict[str, Any]:
    req = urlrequest.Request(
        url,
        method=method,
        data=data,
        headers=headers or {},
    )
    try:
        with urlrequest.urlopen(
            req,
            timeout=timeout,
            context=llm_ssl_context(),
        ) as response:
            raw = response.read().decode("utf-8", errors="replace")
    except urlerror.HTTPError as exc:
        try:
            details = exc.read().decode("utf-8", errors="ignore").strip()
        except Exception:
            details = ""
        message = normalize_provider_error(
            "assemblyai",
            f"HTTP {exc.code}{f': {details}' if details else ''}",
        )
        raise RuntimeError(message) from exc

    try:
        payload = json.loads(raw)
    except json.JSONDecodeError as exc:
        snippet = raw[:500].replace("\n", "\\n")
        raise RuntimeError(f"AssemblyAI вернул невалидный JSON. body={snippet}") from exc
    if not isinstance(payload, dict):
        raise RuntimeError("AssemblyAI вернул невалидный JSON-тип.")
    return payload


def create_assemblyai_transcript(
    config: WorkerConfig,
    audio_url: str,
) -> str:
    payload: dict[str, Any] = {
        "audio_url": audio_url,
        "speaker_labels": True,
        "language_code": config.assemblyai_language_code,
        "speech_models": list(config.assemblyai_speech_models),
    }
    response = assemblyai_request(
        f"{config.assemblyai_base_url}/v2/transcript",
        method="POST",
        data=json.dumps(payload, ensure_ascii=False).encode("utf-8"),
        headers={
            "Authorization": config.assemblyai_api_key,
            "Content-Type": "application/json; charset=utf-8",
            "Accept": "application/json",
        },
        timeout=60.0,
    )
    transcript_id = str(response.get("id") or "").strip()
    if not transcript_id:
        raise RuntimeError("AssemblyAI не вернул id транскрипта.")
    return transcript_id


def get_assemblyai_transcript(
    config: WorkerConfig,
    transcript_id: str,
) -> dict[str, Any]:
    return assemblyai_request(
        f"{config.assemblyai_base_url}/v2/transcript/{transcript_id}",
        headers={
            "Authorization": config.assemblyai_api_key,
            "Accept": "application/json",
        },
        timeout=60.0,
    )


def parse_assemblyai_segments(payload: dict[str, Any]) -> list[dict[str, Any]]:
    utterances = payload.get("utterances")
    if not isinstance(utterances, list):
        raise RuntimeError("AssemblyAI не вернул utterances.")

    speaker_map: dict[str, int] = {}
    segments: list[dict[str, Any]] = []

    for item in utterances:
        if not isinstance(item, dict):
            continue
        text = str(item.get("text") or "").strip()
        start_ms = to_milliseconds(item.get("start"))
        end_ms = to_milliseconds(item.get("end"))
        if not text or start_ms is None or end_ms is None or end_ms < start_ms:
            continue

        speaker_raw = str(item.get("speaker") or "").strip() or "unknown"
        if speaker_raw not in speaker_map:
            speaker_map[speaker_raw] = len(speaker_map)

        segments.append(
            {
                "speaker_id": speaker_map[speaker_raw],
                "text": text,
                "start_ms": start_ms,
                "end_ms": end_ms,
                "confidence": extract_float(item.get("confidence")),
                "provider_payload_json": item,
            }
        )

    if not segments:
        raise RuntimeError("AssemblyAI не вернул diarized utterances.")

    segments.sort(key=lambda item: (item["start_ms"], item["end_ms"]))
    return segments


def transcribe_meeting_with_assemblyai(
    config: WorkerConfig,
    prepared_audio: dict[str, Any],
    source_file: dict[str, Any],
    heartbeat_conn: psycopg.Connection,
    job_id: str,
) -> dict[str, Any]:
    if not config.assemblyai_api_key:
        raise RuntimeError("ASSEMBLYAI_API_KEY не задан.")

    storage_key = str(prepared_audio.get("audio_storage_key") or "").strip()
    if not storage_key:
        storage_key = str(source_file["storage_key"]).strip()

    audio_url = s3_presigned_get_url(
        config,
        storage_key,
        config.assemblyai_audio_url_expires_seconds,
    )
    log_worker_event(
        "assemblyai-transcript-create-start",
        jobId=job_id,
        meetingId=source_file.get("meeting_id"),
        storageKey=storage_key,
    )
    transcript_id = create_assemblyai_transcript(config, audio_url)
    touch_meeting_job_heartbeat(heartbeat_conn, job_id)
    log_worker_event(
        "assemblyai-transcript-created",
        jobId=job_id,
        meetingId=source_file.get("meeting_id"),
        transcriptId=transcript_id,
    )

    started_at = time.time()
    poll_count = 0
    last_status = ""
    while True:
        payload = get_assemblyai_transcript(config, transcript_id)
        touch_meeting_job_heartbeat(heartbeat_conn, job_id)
        poll_count += 1
        status = str(payload.get("status") or "").strip().lower()
        if status != last_status or poll_count == 1 or poll_count % 12 == 0:
            log_worker_event(
                "assemblyai-transcript-status",
                jobId=job_id,
                meetingId=source_file.get("meeting_id"),
                transcriptId=transcript_id,
                status=status,
                pollCount=poll_count,
            )
            last_status = status
        if status == "completed":
            return {
                "transcript_id": transcript_id,
                "audio_url": audio_url,
                "result_payload": payload,
            }
        if status in {"error", "failed"}:
            log_worker_event(
                "assemblyai-transcript-failed",
                jobId=job_id,
                meetingId=source_file.get("meeting_id"),
                transcriptId=transcript_id,
                status=status,
            )
            raise RuntimeError(
                str(
                    payload.get("error")
                    or payload.get("message")
                    or "AssemblyAI transcription failed"
                )
            )
        if time.time() - started_at > config.assemblyai_timeout_seconds:
            raise RuntimeError("AssemblyAI task timed out")
        time.sleep(max(1.0, config.assemblyai_poll_interval_seconds))


def transcribe_meeting_with_salutespeech(
    config: WorkerConfig,
    prepared_audio: dict[str, Any],
    heartbeat_conn: psycopg.Connection,
    job_id: str,
) -> dict[str, Any]:
    token = get_salutespeech_access_token(config)
    request_file_id = salutespeech_upload_file(
        config,
        token,
        str(prepared_audio["file_name"]),
        prepared_audio["audio_bytes"],
    )
    task_id = salutespeech_create_recognition_task(
        config,
        token,
        request_file_id=request_file_id,
        audio_encoding=str(prepared_audio.get("audio_encoding") or "PCM_S16LE"),
        sample_rate=int(prepared_audio["sample_rate"]),
        channels_count=int(prepared_audio["channels_count"]),
    )
    touch_meeting_job_heartbeat(heartbeat_conn, job_id)

    started_at = time.time()
    last_status_payload: dict[str, Any] | None = None
    while True:
        status_payload = salutespeech_get_task_status(config, token, task_id)
        last_status_payload = status_payload
        touch_meeting_job_heartbeat(heartbeat_conn, job_id)
        result = status_payload.get("result") or {}
        status = str(result.get("status") or "").upper()
        if status == "DONE":
            response_file_id = str(result.get("response_file_id") or "").strip()
            if not response_file_id:
                raise RuntimeError("SaluteSpeech вернул DONE без response_file_id.")
            result_payload = salutespeech_download_result(config, token, response_file_id)
            return {
                "task_id": task_id,
                "request_file_id": request_file_id,
                "response_file_id": response_file_id,
                "status_payload": status_payload,
                "result_payload": result_payload,
            }
        if status == "ERROR":
            raise RuntimeError(str(result.get("error") or "SaluteSpeech task failed"))
        if time.time() - started_at > config.salutespeech_timeout_seconds:
            raise RuntimeError("SaluteSpeech task timed out")
        time.sleep(max(1.0, config.salutespeech_poll_interval_seconds))


def transcribe_meeting(
    config: WorkerConfig,
    prepared_audio: dict[str, Any],
    source_file: dict[str, Any],
    heartbeat_conn: psycopg.Connection,
    job_id: str,
) -> dict[str, Any]:
    if config.meeting_transcription_provider == "assemblyai":
        return transcribe_meeting_with_assemblyai(
            config,
            prepared_audio,
            source_file,
            heartbeat_conn,
            job_id,
        )
    return transcribe_meeting_with_salutespeech(
        config,
        prepared_audio,
        heartbeat_conn,
        job_id,
    )


def handle_meeting_audio_prepared(
    conn: psycopg.Connection,
    config: WorkerConfig,
    meeting: dict[str, Any],
    source_file: dict[str, Any],
    heartbeat_conn: psycopg.Connection,
    job_id: str,
) -> dict[str, Any]:
    prepared = prepare_meeting_audio(config, source_file, heartbeat_conn, job_id)
    conn.execute(
        """
        update meeting_source_files
        set
          audio_storage_key = %s,
          audio_mime_type = %s,
          duration_seconds = %s,
          processing_status = 'done'
        where id = %s
        """,
        (
            prepared["audio_storage_key"],
            prepared["audio_mime_type"],
            prepared["duration_seconds"],
            source_file["id"],
        ),
    )
    conn.execute(
        """
        update meetings
        set
          duration_seconds = %s,
          updated_at = now()
        where id = %s
        """,
        (prepared["duration_seconds"], meeting["id"]),
    )
    return {
        "audioStorageKey": prepared["audio_storage_key"],
        "audioMimeType": prepared["audio_mime_type"],
        "durationSeconds": prepared["duration_seconds"],
        "sampleRate": prepared["sample_rate"],
        "channelsCount": prepared["channels_count"],
    }


def handle_meeting_transcript_compiled(
    conn: psycopg.Connection,
    config: WorkerConfig,
    meeting: dict[str, Any],
    source_file: dict[str, Any],
    heartbeat_conn: psycopg.Connection,
    job_id: str,
) -> dict[str, Any]:
    log_worker_event(
        "meeting-transcript-compile-start",
        jobId=job_id,
        meetingId=meeting["id"],
        sourceStorageKey=source_file.get("storage_key"),
        provider=config.meeting_transcription_provider,
    )
    touch_meeting_job_heartbeat(heartbeat_conn, job_id)
    log_worker_event(
        "meeting-audio-prepare-start",
        jobId=job_id,
        meetingId=meeting["id"],
    )
    prepared = prepare_meeting_audio(config, source_file, heartbeat_conn, job_id)
    touch_meeting_job_heartbeat(heartbeat_conn, job_id)
    log_worker_event(
        "meeting-audio-prepare-done",
        jobId=job_id,
        meetingId=meeting["id"],
        durationSeconds=prepared["duration_seconds"],
        channelsCount=prepared["channels_count"],
        sampleRate=prepared["sample_rate"],
        audioStorageKey=prepared["audio_storage_key"],
    )
    log_worker_event(
        "meeting-transcript-provider-start",
        jobId=job_id,
        meetingId=meeting["id"],
        provider=config.meeting_transcription_provider,
    )
    raw_result = transcribe_meeting(
        config,
        prepared,
        source_file,
        heartbeat_conn,
        job_id,
    )
    log_worker_event(
        "meeting-transcript-provider-done",
        jobId=job_id,
        meetingId=meeting["id"],
        provider=config.meeting_transcription_provider,
        resultKeys=sorted(raw_result.keys()),
    )
    if config.meeting_transcription_provider == "assemblyai":
        diarized_segments = parse_assemblyai_segments(raw_result["result_payload"])
    else:
        diarized_segments = parse_salutespeech_segments(raw_result["result_payload"])
    speakers, segments = upsert_meeting_transcript(conn, meeting["id"], diarized_segments)
    transcript_markdown = build_meeting_transcript_markdown(
        str(meeting["title"]),
        segments,
    )
    transcript_json = build_meeting_transcript_json(meeting["id"], speakers, segments)

    ensure_meeting_artifact(
        conn,
        meeting["id"],
        "transcript_compiled",
        transcript_markdown,
        transcript_json,
    )
    conn.execute(
        """
        update meeting_source_files
        set
          audio_storage_key = %s,
          audio_mime_type = %s,
          duration_seconds = %s,
          processing_status = 'done'
        where id = %s
        """,
        (
            prepared["audio_storage_key"],
            prepared["audio_mime_type"],
            prepared["duration_seconds"],
            source_file["id"],
        ),
    )
    conn.execute(
        """
        update meetings
        set
          duration_seconds = %s,
          speakers_count = %s,
          updated_at = now()
        where id = %s
        """,
        (prepared["duration_seconds"], len(speakers), meeting["id"]),
    )
    return {
        "provider": config.meeting_transcription_provider,
        "model": (
            list(config.assemblyai_speech_models)[0]
            if config.meeting_transcription_provider == "assemblyai"
            else config.salutespeech_model
        ),
        "language": (
            config.assemblyai_language_code
            if config.meeting_transcription_provider == "assemblyai"
            else config.salutespeech_language
        ),
        "speakerSeparation": True,
        "audioStorageKey": prepared["audio_storage_key"],
        "sampleRate": prepared["sample_rate"],
        "channelsCount": prepared["channels_count"],
        "durationSeconds": prepared["duration_seconds"],
        "speakersCount": len(speakers),
        "segmentsCount": len(segments),
        **{
            config.meeting_transcription_provider: raw_result,
        },
    }


def handle_meeting_generation_stage(
    conn: psycopg.Connection,
    config: WorkerConfig,
    meeting: dict[str, Any],
    stage: str,
    heartbeat_conn: psycopg.Connection,
    job_id: str,
) -> dict[str, Any]:
    transcript_json_artifact = load_meeting_artifact(conn, meeting["id"], "transcript_compiled", "json")
    if transcript_json_artifact is None:
        raise RuntimeError("Transcript artifact not found for meeting generation.")
    transcript_payload = transcript_json_artifact.get("content_json") or {}
    transcript_data = parse_json(transcript_payload, {})
    speakers = transcript_data.get("speakers")
    segments = transcript_data.get("segments")
    if not isinstance(segments, list) or not segments:
        raise RuntimeError("Transcript JSON не содержит segments.")

    normalized_segments: list[dict[str, Any]] = []
    speakers_by_id: dict[str, dict[str, Any]] = {}
    if isinstance(speakers, list):
        for speaker in speakers:
            if isinstance(speaker, dict):
                speakers_by_id[str(speaker.get("id") or "")] = speaker

    for segment in segments:
        if not isinstance(segment, dict):
            continue
        speaker_id = str(segment.get("speakerId") or "")
        display_name = str(segment.get("displayName") or "").strip()
        if not display_name and speaker_id:
            display_name = str(speakers_by_id.get(speaker_id, {}).get("displayName") or "").strip()
        normalized_segments.append(
            {
                "id": int(segment.get("id") or 0),
                "speaker_id": speaker_id or None,
                "speaker_label": str(segment.get("speakerLabel") or ""),
                "display_name": display_name or "Спикер",
                "start_ms": int(segment.get("startMs") or 0),
                "end_ms": int(segment.get("endMs") or 0),
                "text": str(segment.get("text") or "").strip(),
                "confidence": segment.get("confidence"),
            }
        )

    transcript_input = build_meeting_llm_transcript_input(normalized_segments)
    touch_meeting_job_heartbeat(heartbeat_conn, job_id)
    llm_result = generate_meeting_markdown_with_llm(
        conn,
        config,
        stage=stage,
        meeting_title=str(meeting["title"]),
        transcript_input=transcript_input,
    )
    touch_meeting_job_heartbeat(heartbeat_conn, job_id)
    markdown = str(llm_result.get("markdown") or "").strip()
    if stage == "meeting_actions":
        content_json = meeting_actions_json_from_llm(llm_result, markdown)
    else:
        content_json = {
            "stage": stage,
            "markdown": markdown,
            "shortSummary": str(llm_result.get("shortSummary") or "").strip(),
        }

    ensure_meeting_artifact(conn, meeting["id"], stage, markdown, content_json)
    return {
        "llm": {
            "promptModule": llm_result.get("promptModule"),
            "promptKey": llm_result.get("promptKey"),
            "promptTitle": llm_result.get("promptTitle"),
            "provider": llm_result.get("provider"),
            "model": llm_result.get("model"),
            "attempts": llm_result.get("attempts"),
        },
        "markdownLength": len(markdown),
        "transcriptSegments": len(normalized_segments),
    }


def queue_next_meeting_job(
    conn: psycopg.Connection,
    config: WorkerConfig,
    meeting_id: str,
    stage: str,
) -> dict[str, Any] | None:
    next_stage_value = next_meeting_stage(stage)
    if not next_stage_value:
        return None
    next_job_id = f"{meeting_id}-{next_stage_value}-{int(time.time() * 1000)}"
    conn.execute(
        """
        insert into meeting_jobs (
          id, meeting_id, stage, status, payload_json, result_json, error_text, started_at, finished_at, created_at
        ) values (
          %s, %s, %s, 'queued', %s::jsonb, null, null, null, null, now()
        )
        """,
        (
            next_job_id,
            meeting_id,
            next_stage_value,
            json.dumps({"stage": next_stage_value, "trigger": "auto"}, ensure_ascii=False),
        ),
    )
    return {
        "jobId": next_job_id,
        "meetingId": meeting_id,
        "stage": next_stage_value,
        "trigger": "auto",
    }


def process_meeting_job(config: WorkerConfig, job_message: dict[str, Any]) -> None:
    conn = psycopg.connect(config.postgres_url, row_factory=dict_row, autocommit=True)
    status_conn = psycopg.connect(config.postgres_url, row_factory=dict_row, autocommit=True)
    try:
        job = load_meeting_job(conn, job_message["jobId"])
        if job is None:
            print(
                json.dumps(
                    {
                        "service": "worker",
                        "status": "skipped",
                        "reason": "meeting-job-missing",
                        "jobId": job_message["jobId"],
                    },
                    ensure_ascii=False,
                )
            )
            return

        if job["status"] == "done":
            print(
                json.dumps(
                    {
                        "service": "worker",
                        "status": "skipped",
                        "reason": "meeting-job-already-done",
                        "jobId": job_message["jobId"],
                    },
                    ensure_ascii=False,
                )
            )
            return

        meeting = load_meeting(conn, job["meeting_id"])
        source_file = load_meeting_source_file(conn, meeting["id"])
        if source_file is None:
            raise RuntimeError("Meeting source file not found.")

        metadata: dict[str, Any]
        next_job_payload: dict[str, Any] | None = None

        set_meeting_job_processing(status_conn, job["id"])
        with conn.transaction():
            conn.execute(
                """
                update meeting_source_files
                set processing_status = 'processing'
                where id = %s
                """,
                (source_file["id"],),
            )
            conn.execute(
                """
                update meetings
                set
                  status = 'processing',
                  processing_started_at = coalesce(processing_started_at, now()),
                  processing_finished_at = null,
                  error_text = null,
                  updated_at = now()
                where id = %s
                """,
                (meeting["id"],),
            )

            if job["stage"] == "audio_prepared":
                metadata = handle_meeting_audio_prepared(
                    conn,
                    config,
                    meeting,
                    source_file,
                    status_conn,
                    job["id"],
                )
            elif job["stage"] == "transcript_compiled":
                metadata = handle_meeting_transcript_compiled(
                    conn,
                    config,
                    meeting,
                    source_file,
                    status_conn,
                    job["id"],
                )
            else:
                metadata = handle_meeting_generation_stage(
                    conn,
                    config,
                    meeting,
                    job["stage"],
                    status_conn,
                    job["id"],
                )

            trigger = str((job.get("payload_json") or {}).get("trigger") or "")
            if trigger in {"start", "auto"}:
                next_job_payload = queue_next_meeting_job(
                    conn, config, meeting["id"], job["stage"]
                )

            if next_job_payload is None:
                conn.execute(
                    """
                    update meetings
                    set
                      status = 'completed',
                      processing_finished_at = now(),
                      updated_at = now()
                    where id = %s
                    """,
                    (meeting["id"],),
                )

        set_meeting_job_done(status_conn, job["id"], job["stage"], metadata)

        if next_job_payload is not None:
            redis_command(
                config.redis_url,
                "LPUSH",
                config.meeting_job_queue_key,
                json.dumps(next_job_payload, ensure_ascii=False),
            )

        print(
            json.dumps(
                {
                    "service": "worker",
                    "status": "done",
                    "jobType": "meeting",
                    "jobId": job["id"],
                    "meetingId": meeting["id"],
                    "stage": job["stage"],
                },
                ensure_ascii=False,
            )
        )
    except Exception as exc:
        try:
            conn.rollback()
        except Exception:
            pass
        try:
            set_meeting_job_failed(status_conn, job_message["jobId"], str(exc))
            conn.execute(
                """
                update meetings
                set
                  status = 'failed',
                  error_text = %s,
                  updated_at = now()
                where id = %s
                """,
                (str(exc), job_message["meetingId"]),
            )
            conn.execute(
                """
                update meeting_source_files
                set processing_status = 'failed'
                where meeting_id = %s
                """,
                (job_message["meetingId"],),
            )
        except Exception:
            pass
        print(
            json.dumps(
                {
                    "service": "worker",
                    "status": "failed",
                    "jobType": "meeting",
                    "error": str(exc),
                },
                ensure_ascii=False,
            )
        )
    finally:
        try:
            status_conn.close()
        except Exception:
            pass
        conn.close()


def process_job(config: WorkerConfig, job_message: dict[str, Any]) -> None:
    # Use autocommit mode so `with conn.transaction()` always creates a real
    # top-level transaction. Otherwise prior SELECT opens an implicit outer
    # transaction and later writes can be rolled back on connection close.
    conn = psycopg.connect(config.postgres_url, row_factory=dict_row, autocommit=True)
    try:
        job = load_job(conn, job_message["jobId"])
        if job is None:
            try:
                with conn.transaction():
                    recovered = recover_missing_job(conn, job_message)
                if recovered is not None:
                    job = recovered
                    print(
                        json.dumps(
                            {
                                "service": "worker",
                                "status": "recovered",
                                "reason": "job-created-from-queue-message",
                                "jobId": job_message["jobId"],
                            },
                            ensure_ascii=False,
                        )
                    )
            except Exception:
                recovered = None

        if job is None:
            print(
                json.dumps(
                    {
                        "service": "worker",
                        "status": "skipped",
                        "reason": "job-missing",
                        "jobId": job_message["jobId"],
                    },
                    ensure_ascii=False,
                )
            )
            return

        if job["status"] == "done":
            print(
                json.dumps(
                    {
                        "service": "worker",
                        "status": "skipped",
                        "reason": "already-done",
                        "jobId": job_message["jobId"],
                    },
                    ensure_ascii=False,
                )
            )
            return

        project = load_project(conn, job["project_id"])
        transcription_data: dict[str, Any] | None = None
        llm_data: dict[str, Any] | None = None
        llm_stage_data: dict[str, Any] | None = None
        llm_documents: list[dict[str, Any]] = []
        llm_stage_input_stages: list[str] = []
        source_file: dict[str, Any] | None = None
        source_compiled_sections: list[str] = []
        topic_for_markdown = project["source_summary"]
        if job["stage"] == "source_compiled":
            source_files = load_source_files(conn, project["id"])
            source_file = pick_video_source_file(source_files)
            document_source_files = pick_document_source_files(source_files)
            if source_file is not None:
                transcription_data = call_transcription_service(config, source_file)
                transcript_text = str(transcription_data.get("text") or "").strip()
                llm_data = analyze_source_with_llm(
                    conn,
                    config,
                    transcript_text,
                    prompt_key="analize_video",
                    source_type="транскрипт видео",
                )
                video_source_text = str(llm_data.get("transcript") or "").strip() or transcript_text
                if video_source_text:
                    source_compiled_sections.append(
                        "\n".join(
                            [
                                f"## Видео: {source_file.get('original_name') or source_file['id']}",
                                video_source_text,
                            ]
                        )
                    )
                video_summary = str(llm_data.get("summary") or "").strip() or extract_summary(
                    video_source_text
                )
                if video_summary:
                    topic_for_markdown = video_summary
            if document_source_files:
                merged_document_parts: list[str] = []
                for document_source in document_source_files:
                    try:
                        document_bytes = download_s3_object_bytes(
                            config, str(document_source["storage_key"])
                        )
                        document_text = extract_document_text(document_source, document_bytes)
                        document_text = document_text.strip()
                        if not document_text:
                            continue
                        source_compiled_sections.append(
                            "\n".join(
                                [
                                    f"## Документ: {document_source.get('original_name') or document_source['id']}",
                                    document_text,
                                ]
                            )
                        )
                        merged_document_parts.append(
                            "\n\n".join(
                                [
                                    f"Файл: {document_source.get('original_name') or document_source['id']}",
                                    document_text[:50000],
                                ]
                            )
                        )
                    except Exception as exc:
                        llm_documents.append(
                            {
                                "fileId": document_source["id"],
                                "sourceOriginalName": document_source.get("original_name"),
                                "error": str(exc),
                            }
                        )
                if merged_document_parts:
                    merged_document_text = "\n\n---\n\n".join(merged_document_parts)[:180000]
                    llm_document_data = analyze_source_with_llm(
                        conn,
                        config,
                        merged_document_text,
                        prompt_key="analize_doc",
                        source_type="структурированные документы",
                    )
                    llm_documents.append(
                        {
                            "promptModule": llm_document_data.get("promptModule"),
                            "promptKey": llm_document_data.get("promptKey"),
                            "promptTitle": llm_document_data.get("promptTitle"),
                            "provider": llm_document_data.get("provider"),
                            "model": llm_document_data.get("model"),
                            "attempts": llm_document_data.get("attempts"),
                            "summaryLength": len(str(llm_document_data.get("summary") or "")),
                            "sourceTextLength": len(
                                str(llm_document_data.get("transcript") or "")
                            ),
                            "files": [
                                {
                                    "fileId": item["id"],
                                    "sourceStorageKey": item.get("storage_key"),
                                    "sourceOriginalName": item.get("original_name"),
                                    "sourceMimeType": item.get("mime_type"),
                                }
                                for item in document_source_files
                            ],
                        }
                    )
                    document_summary = str(llm_document_data.get("summary") or "").strip()
                    if document_summary:
                        topic_for_markdown = (
                            f"{topic_for_markdown}\n\n{document_summary}"
                            if topic_for_markdown.strip()
                            else document_summary
                        )
        else:
            source_for_stage, llm_stage_input_stages = build_stage_generation_input(
                conn,
                project["id"],
                job["stage"],
                str(project.get("source_summary") or ""),
            )
            if source_for_stage.strip():
                llm_stage_data = generate_stage_markdown_with_llm(
                    conn,
                    config,
                    stage=job["stage"],
                    source_text=source_for_stage,
                    project_name=str(project["name"]),
                )
                topic_for_markdown = (
                    str(llm_stage_data.get("shortSummary") or "").strip()
                    or extract_summary(source_for_stage)
                )

        markdown = str(llm_stage_data.get("markdown") or "").strip() if llm_stage_data else ""
        if job["stage"] == "source_compiled":
            markdown = make_source_compiled_markdown(project["name"], source_compiled_sections)
        if not markdown:
            markdown = make_stage_markdown(project["name"], job["stage"], topic_for_markdown)
        payload = job.get("payload_json") or {}
        if isinstance(payload, str):
            payload = json.loads(payload)
        prompt_metadata = stage_prompt_metadata(conn, job["stage"])
        stage_metadata = dict(prompt_metadata)
        if transcription_data is not None and source_file is not None:
            stage_metadata["transcription"] = {
                "sourceFileId": source_file["id"],
                "sourceStorageKey": source_file["storage_key"],
                "sourceOriginalName": source_file.get("original_name"),
                "sourceMimeType": source_file.get("mime_type"),
                "model": transcription_data.get("model"),
                "duration": transcription_data.get("duration"),
                "textLength": len(str(transcription_data.get("text") or "")),
                "segmentsCount": len(transcription_data.get("segments") or []),
            }
        if llm_data is not None:
            stage_metadata["llm"] = {
                "promptModule": llm_data.get("promptModule"),
                "promptKey": llm_data.get("promptKey"),
                "promptTitle": llm_data.get("promptTitle"),
                "provider": llm_data.get("provider"),
                "model": llm_data.get("model"),
                "attempts": llm_data.get("attempts"),
                "summaryLength": len(str(llm_data.get("summary") or "")),
                "sourceTextLength": len(str(llm_data.get("transcript") or "")),
            }
        if llm_documents:
            stage_metadata["llmDocuments"] = llm_documents
        if llm_stage_data is not None:
            stage_metadata["llmStageGeneration"] = {
                "promptModule": llm_stage_data.get("promptModule"),
                "provider": llm_stage_data.get("provider"),
                "model": llm_stage_data.get("model"),
                "attempts": llm_stage_data.get("attempts"),
                "promptKey": llm_stage_data.get("promptKey"),
                "promptTitle": llm_stage_data.get("promptTitle"),
                "sourceTextLength": llm_stage_data.get("sourceTextLength"),
                "inputStages": llm_stage_input_stages,
                "markdownLength": len(str(llm_stage_data.get("markdown") or "")),
                "shortSummaryLength": len(str(llm_stage_data.get("shortSummary") or "")),
            }
        auto_generate_all = bool(payload.get("autoGenerateAll"))
        queued_next_stage = payload.get("nextStage")
        if queued_next_stage is not None and not isinstance(queued_next_stage, str):
            queued_next_stage = None
        resolved_next_stage = queued_next_stage or (
            next_stage(job["stage"]) if auto_generate_all else None
        )
        should_finish_course = resolved_next_stage is None and job["stage"] == "course_test"
        next_job_payload = None

        with conn.transaction():
            set_job_processing(conn, job["id"])
            append_project_log(conn, project["id"], f"Worker начал обработку job {job['stage']}.")
            ensure_stage_artifact(conn, job, markdown, stage_metadata)
            set_job_done(conn, job["id"], job["stage"], stage_metadata)
            if resolved_next_stage:
                next_job_id = f"{project['id']}-{resolved_next_stage}-{int(time.time() * 1000)}"
                next_job_payload = {
                    "jobId": next_job_id,
                    "projectId": project["id"],
                    "stage": resolved_next_stage,
                    "trigger": "auto" if auto_generate_all else "manual",
                }
                conn.execute(
                    """
                    insert into processing_jobs (
                      id, project_id, stage, status, payload_json, result_json, error_text, started_at, finished_at, created_at
                    ) values (
                      %s, %s, %s, 'queued', %s::jsonb, null, null, null, null, now()
                    )
                    """,
                    (
                        next_job_id,
                        project["id"],
                        resolved_next_stage,
                        json.dumps(
                            {
                                "stage": resolved_next_stage,
                                "trigger": "auto" if auto_generate_all else "manual",
                                "autoGenerateAll": auto_generate_all,
                                "nextStage": next_stage(resolved_next_stage)
                                if auto_generate_all
                                else None,
                            },
                            ensure_ascii=False,
                        ),
                    ),
                )
                conn.execute(
                    """
                    update projects
                    set
                      current_stage = %s,
                      status = 'processing',
                      progress = %s,
                      updated_at = now()
                    where id = %s
                    """,
                    (
                        resolved_next_stage,
                        progress_for_stage(resolved_next_stage, completed=False),
                        project["id"],
                    ),
                )
            else:
                conn.execute(
                    """
                    update projects
                    set
                      current_stage = %s,
                      status = %s,
                      progress = %s,
                      updated_at = now()
                    where id = %s
                    """,
                    (
                        job["stage"],
                        "completed" if should_finish_course else "uploaded",
                        progress_for_stage(job["stage"], completed=should_finish_course),
                        project["id"],
                    ),
                )

        if next_job_payload:
            redis_command(
                config.redis_url,
                "LPUSH",
                config.job_queue_key,
                json.dumps(next_job_payload, ensure_ascii=False),
            )

        print(
            json.dumps(
                {
                    "service": "worker",
                    "status": "done",
                    "jobId": job["id"],
                    "projectId": project["id"],
                    "stage": job["stage"],
                },
                ensure_ascii=False,
            )
        )
    except Exception as exc:
        try:
            conn.rollback()
        except Exception:
            pass
        try:
            set_job_failed(conn, job_message["jobId"], str(exc))
            conn.execute(
                """
                update projects
                set status = 'failed', updated_at = now()
                where id = %s
                """,
                (job_message["projectId"],),
            )
            conn.commit()
        except Exception:
            pass
        print(
            json.dumps(
                {"service": "worker", "status": "failed", "error": str(exc)},
                ensure_ascii=False,
            )
        )
    finally:
        conn.close()


def drain_queue(config: WorkerConfig) -> None:
    print(
        json.dumps(
            {
                "service": "worker",
                "config": asdict(config),
                "queues": [config.job_queue_key, config.meeting_job_queue_key],
            },
            ensure_ascii=False,
            indent=2,
        )
    )

    while True:
        try:
            raw_message = redis_command(
                config.redis_url,
                "BRPOP",
                config.job_queue_key,
                config.meeting_job_queue_key,
                "5",
                timeout=6,
            )
            if not raw_message:
                print(json.dumps({"service": "worker", "status": "heartbeat"}, ensure_ascii=False))
                continue

            queue_name = config.job_queue_key
            if isinstance(raw_message, list):
                queue_name = str(raw_message[0])
                raw_message = raw_message[-1]

            message = json.loads(raw_message)
            if queue_name == config.meeting_job_queue_key:
                process_meeting_job(config, message)
            else:
                process_job(config, message)
        except KeyboardInterrupt:
            print(json.dumps({"service": "worker", "status": "stopped"}, ensure_ascii=False))
            return
        except Exception as exc:
            print(
                json.dumps(
                    {"service": "worker", "status": "error", "error": str(exc)},
                    ensure_ascii=False,
                )
            )
            time.sleep(2)


def main() -> None:
    drain_queue(load_config())


if __name__ == "__main__":
    main()
